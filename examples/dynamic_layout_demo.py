"""
Dynamic Layout Demo - Create Custom Slide Layouts Without Template Constraints

This example demonstrates how PPTX Agent creates flexible, LLM-driven layouts:
- Templates provide STYLING (colors, fonts) not STRUCTURE
- Layouts are calculated programmatically based on content
- LLM can make layout decisions dynamically
- Custom arrangements not defined in templates are possible

Key Benefits:
- Not locked into predefined template layouts
- Mix and match elements creatively
- Adapt layout to content requirements
- Use templates for visual consistency, not structural constraints
"""

from pathlib import Path
from pptx_agent.core.pptx_handler import PPTXHandler
from pptx_agent.core.dynamic_layout_engine import DynamicLayoutEngine, LayoutRegion, ElementBox
from pptx_agent.core.template_manager import TemplateManager
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def example_1_basic_dynamic_layout():
    """
    Example 1: Basic dynamic layout without template constraints.

    Creates slides using algorithmic positioning, not template layouts.
    """
    print("="*70)
    print("Example 1: Basic Dynamic Layout")
    print("="*70)

    # Initialize handler (no template)
    handler = PPTXHandler()

    # Initialize dynamic layout engine
    layout_engine = DynamicLayoutEngine(
        slide_width=10.0,   # Standard 4:3
        slide_height=7.5
    )

    # Create a blank slide (layout 6 is typically blank)
    slide = handler.add_slide(6)

    # Get a simple title + content layout
    elements = layout_engine.create_title_content_layout()

    # Add title
    title_box = elements[0]
    title = handler.add_text_box(
        slide,
        "Dynamically Positioned Title",
        left=title_box.left,
        top=title_box.top,
        width=title_box.width,
        height=title_box.height,
        font_size=36,
        bold=True
    )

    # Add content
    content_box = elements[1]
    content = handler.add_text_box(
        slide,
        "This content area was positioned algorithmically by the DynamicLayoutEngine,\n"
        "not by a predefined template layout.\n\n"
        "Benefits:\n"
        "• Not constrained by template structure\n"
        "• Can create custom arrangements\n"
        "• LLM can make layout decisions\n"
        "• Templates only provide colors/fonts",
        left=content_box.left,
        top=content_box.top,
        width=content_box.width,
        height=content_box.height,
        font_size=18
    )

    handler.save(Path("dynamic_basic.pptx"))
    print("✓ Created: dynamic_basic.pptx")
    print("  - Title + content layout calculated programmatically")
    print("  - No template layout used\n")


def example_2_two_column_layout():
    """
    Example 2: Dynamic two-column layout with custom ratio.

    LLM can specify column ratios based on content needs.
    """
    print("="*70)
    print("Example 2: Dynamic Two-Column Layout")
    print("="*70)

    handler = PPTXHandler()
    layout_engine = DynamicLayoutEngine(slide_width=13.333, slide_height=7.5)  # Widescreen

    slide = handler.add_slide(6)  # Blank

    # Create 60/40 split (LLM could decide this ratio)
    elements = layout_engine.create_two_column_layout(
        title=True,
        column_ratio=0.6  # 60% left, 40% right
    )

    # Add title
    title_elem = [e for e in elements if e.element_type == 'title'][0]
    handler.add_text_box(
        slide, "Custom Column Ratio: 60/40",
        title_elem.left, title_elem.top, title_elem.width, title_elem.height,
        font_size=32, bold=True
    )

    # Left column (larger)
    left_elem = [e for e in elements if e.element_type == 'column_left'][0]
    handler.add_text_box(
        slide,
        "Main Content (60%)\n\n"
        "This column gets more space because the LLM determined the main content "
        "needs more room.\n\n"
        "Template layouts are often 50/50, but we can customize based on actual content needs.",
        left_elem.left, left_elem.top, left_elem.width, left_elem.height,
        font_size=16
    )

    # Right column (smaller)
    right_elem = [e for e in elements if e.element_type == 'column_right'][0]
    handler.add_text_box(
        slide,
        "Sidebar (40%)\n\n"
        "Supporting information or notes.",
        right_elem.left, right_elem.top, right_elem.width, right_elem.height,
        font_size=14
    )

    handler.save(Path("dynamic_columns.pptx"))
    print("✓ Created: dynamic_columns.pptx")
    print("  - Custom 60/40 column split")
    print("  - LLM-driven ratio decision\n")


def example_3_grid_layout():
    """
    Example 3: Dynamic grid layout for images.

    Grid dimensions calculated based on number of items.
    """
    print("="*70)
    print("Example 3: Dynamic Grid Layout")
    print("="*70)

    handler = PPTXHandler()
    layout_engine = DynamicLayoutEngine(slide_width=13.333, slide_height=7.5)

    slide = handler.add_slide(6)

    # Create 2x3 grid
    elements = layout_engine.create_grid_layout(rows=2, cols=3, title=True)

    # Title
    title_elem = [e for e in elements if e.element_type == 'title'][0]
    handler.add_text_box(
        slide, "Dynamic 2x3 Grid Layout",
        title_elem.left, title_elem.top, title_elem.width, title_elem.height,
        font_size=32, bold=True
    )

    # Fill grid cells
    grid_cells = [e for e in elements if 'grid_cell' in e.element_type]
    colors = [
        (255, 182, 193), (173, 216, 230), (255, 218, 185),
        (221, 160, 221), (144, 238, 144), (255, 228, 196)
    ]

    for i, cell in enumerate(grid_cells):
        # Add colored rectangle
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cell.left), Inches(cell.top),
            Inches(cell.width), Inches(cell.height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*colors[i])
        shape.text = f"Cell {i+1}"
        shape.text_frame.paragraphs[0].alignment = 1  # Center

    handler.save(Path("dynamic_grid.pptx"))
    print("✓ Created: dynamic_grid.pptx")
    print("  - 2x3 grid calculated dynamically")
    print("  - Perfect for image galleries\n")


def example_4_custom_llm_layout():
    """
    Example 4: Completely custom layout from LLM specifications.

    LLM provides exact positioning - total freedom.
    """
    print("="*70)
    print("Example 4: LLM-Specified Custom Layout")
    print("="*70)

    handler = PPTXHandler()
    layout_engine = DynamicLayoutEngine(slide_width=10.0, slide_height=7.5)

    slide = handler.add_slide(6)

    # LLM could generate this specification
    llm_layout_spec = [
        {
            'type': 'title',
            'left': 0.5,
            'top': 0.5,
            'width': 9.0,
            'height': 0.8
        },
        {
            'type': 'main_chart',
            'left': 0.5,
            'top': 1.5,
            'width': 5.5,
            'height': 4.0
        },
        {
            'type': 'key_takeaways',
            'left': 6.2,
            'top': 1.5,
            'width': 3.3,
            'height': 4.0
        },
        {
            'type': 'footer_note',
            'left': 0.5,
            'top': 5.8,
            'width': 9.0,
            'height': 0.5
        }
    ]

    # Create layout from LLM specs
    elements = layout_engine.create_custom_layout(llm_layout_spec)

    # Render each element
    for elem in elements:
        if elem.element_type == 'title':
            handler.add_text_box(
                slide, "LLM-Designed Custom Layout",
                elem.left, elem.top, elem.width, elem.height,
                font_size=32, bold=True
            )
        elif elem.element_type == 'main_chart':
            # Add placeholder for chart
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(elem.left), Inches(elem.top),
                Inches(elem.width), Inches(elem.height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200, 220, 240)
            shape.text = "Chart Area\n(LLM positioned)"
        elif elem.element_type == 'key_takeaways':
            handler.add_text_box(
                slide,
                "Key Takeaways:\n\n"
                "• LLM specified exact positions\n"
                "• Not bound by templates\n"
                "• Custom arrangements\n"
                "• Content-driven layout",
                elem.left, elem.top, elem.width, elem.height,
                font_size=14
            )
        elif elem.element_type == 'footer_note':
            handler.add_text_box(
                slide, "Source: AI-generated layout specification",
                elem.left, elem.top, elem.width, elem.height,
                font_size=10
            )

    handler.save(Path("dynamic_llm_custom.pptx"))
    print("✓ Created: dynamic_llm_custom.pptx")
    print("  - Completely custom layout from LLM")
    print("  - No template constraints\n")


def example_5_template_styling_only():
    """
    Example 5: Use template for styling, not structure.

    Extract colors/fonts from template, but create custom layouts.
    """
    print("="*70)
    print("Example 5: Template Styling + Dynamic Layout")
    print("="*70)

    # Note: This demonstrates the concept even without an actual template file
    template_path = Path("corporate_template.pptx")

    if template_path.exists():
        # Load template for styling only
        template_mgr = TemplateManager(template_path)
        theme = template_mgr.get_theme_summary()

        print(f"Loaded template: {template_path.name}")
        print(f"Extracted theme:")
        print(f"  Colors: {len(theme['colors'])} theme colors")
        print(f"  Fonts: {theme['fonts']}")

        # Use template colors in custom layout
        primary_color = template_mgr.get_theme_color('primary')
        title_style = template_mgr.get_text_style('title')
    else:
        print(f"Note: {template_path} not found (this is a demonstration)")
        print("In practice:")
        print("  1. TemplateManager extracts colors/fonts from template")
        print("  2. DynamicLayoutEngine creates custom positioning")
        print("  3. Elements use template colors but custom positions")
        print("\nKey Concept:")
        print("  Template = Styling (colors, fonts)")
        print("  Layout = Dynamic (positions, sizes)")
        primary_color = (68, 114, 196)
        title_style = {'font_size': 36, 'bold': True}

    # Create custom layout with template styling
    handler = PPTXHandler()
    layout_engine = DynamicLayoutEngine()
    slide = handler.add_slide(6)

    # Custom layout
    elements = layout_engine.create_two_column_layout(title=True, column_ratio=0.65)

    # Title with template styling
    title_elem = [e for e in elements if e.element_type == 'title'][0]
    handler.add_text_box(
        slide, "Template Colors + Custom Layout",
        title_elem.left, title_elem.top, title_elem.width, title_elem.height,
        font_size=title_style['font_size'],
        bold=title_style['bold'],
        color=primary_color
    )

    handler.save(Path("dynamic_template_styling.pptx"))
    print("\n✓ Created: dynamic_template_styling.pptx")
    print("  - Used template for colors/fonts")
    print("  - Used DynamicLayoutEngine for positioning\n")


def example_6_content_aware_layout():
    """
    Example 6: Layout adapts to content requirements.

    DynamicLayoutEngine suggests optimal layout based on content.
    """
    print("="*70)
    print("Example 6: Content-Aware Layout Suggestions")
    print("="*70)

    handler = PPTXHandler()
    layout_engine = DynamicLayoutEngine(slide_width=13.333, slide_height=7.5)

    # Scenario 1: Content with image
    content_desc_1 = {
        'has_title': True,
        'has_image': True,
        'has_chart': False,
        'comparison': False
    }
    suggested_layout_1 = layout_engine.suggest_layout(content_desc_1)
    print(f"\nContent: Title + Image")
    print(f"Suggested: {len(suggested_layout_1)} elements - {[e.element_type for e in suggested_layout_1]}")

    # Scenario 2: Comparison content
    content_desc_2 = {
        'has_title': True,
        'comparison': True
    }
    suggested_layout_2 = layout_engine.suggest_layout(content_desc_2)
    print(f"\nContent: Comparison")
    print(f"Suggested: {len(suggested_layout_2)} elements - {[e.element_type for e in suggested_layout_2]}")

    # Scenario 3: Multiple images
    content_desc_3 = {
        'has_title': True,
        'num_images': 4
    }
    suggested_layout_3 = layout_engine.suggest_layout(content_desc_3)
    print(f"\nContent: 4 images")
    print(f"Suggested: {len(suggested_layout_3)} elements - 2x2 grid")

    print("\n✓ Layout engine adapts to content requirements")
    print("  - Not fixed templates")
    print("  - Dynamic decisions\n")


def main():
    """Run all examples."""

    print("\n" + "="*70)
    print(" Dynamic Layout Engine - Freedom from Template Constraints")
    print("="*70 + "\n")

    print("This demonstration shows how PPTX Agent creates flexible layouts:\n")
    print("Key Principles:")
    print("  1. Templates provide STYLING (colors, fonts), not STRUCTURE")
    print("  2. Layouts are calculated DYNAMICALLY based on content")
    print("  3. LLM can make LAYOUT DECISIONS, not just content")
    print("  4. CUSTOM arrangements not defined in templates are possible\n")

    # Run examples
    example_1_basic_dynamic_layout()
    example_2_two_column_layout()
    example_3_grid_layout()
    example_4_custom_llm_layout()
    example_5_template_styling_only()
    example_6_content_aware_layout()

    print("="*70)
    print("Summary")
    print("="*70)
    print("""
The Dynamic Layout Engine enables:

✓ LLM-Driven Layout Decisions
  - AI can specify custom positioning
  - Not locked into template layouts
  - Adapt to content requirements

✓ Template Flexibility
  - Extract colors/fonts from templates
  - Create custom arrangements
  - Visual consistency + structural freedom

✓ Content-Aware Positioning
  - Grid layouts for multiple images
  - Column ratios based on content
  - Optimal arrangements

✓ Programmatic Control
  - Algorithmic positioning
  - Precise element placement
  - Custom specifications

This makes PPTX Agent truly flexible - templates are a starting point
for styling, not a constraint on structure.
    """)


if __name__ == "__main__":
    main()
