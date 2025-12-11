"""
Flexible Template Usage Demo - Preserve Branding While Customizing Content

This example demonstrates the PPTX Agent's ability to:
1. Understand template slides (identify branding vs content)
2. Preserve branding elements (logos, watermarks, headers/footers)
3. Flexibly position content within safe areas
4. Start from template slides, not blank canvases

Key Concept:
-----------
Traditional approach: Locked into template layout structure
PPTX Agent approach: Extract branding + dynamic content positioning

Benefits:
---------
✓ Keep brand consistency (logos, colors, fonts)
✓ Flexible content layouts not defined in template
✓ LLM can make layout decisions
✓ Safe positioning that doesn't overlap branding
"""

from pathlib import Path
from pptx_agent.core.flexible_slide_builder import FlexibleSlideBuilder
from pptx_agent.core.template_slide_parser import TemplateSlideParser


def example_1_analyze_template():
    """
    Example 1: Analyze a template to understand its structure.

    This shows how PPTX Agent parses templates to identify branding.
    """
    print("="*70)
    print("Example 1: Template Analysis")
    print("="*70)

    # Note: In practice, you would have an actual template file
    template_path = Path("corporate_template.pptx")

    if template_path.exists():
        parser = TemplateSlideParser(template_path)

        # Analyze first few layouts
        print("\nAnalyzing template layouts...\n")

        for layout_idx in range(min(3, len(parser.presentation.slide_layouts))):
            print(parser.get_layout_summary(layout_idx))
            print("\n" + "-"*70 + "\n")

        print("✓ Template analyzed successfully")
        print("  - Identified branding elements (logos, headers, footers)")
        print("  - Calculated safe content areas")
        print("  - Ready for flexible content positioning\n")

    else:
        print(f"Note: {template_path} not found. This is a demonstration.")
        print("\nWhat the analyzer would show:")
        print("------------------------------")
        print("Layout: Title Slide (Index 0)")
        print("Dimensions: 13.3\" x 7.5\"")
        print("")
        print("Total Elements: 5")
        print("  Branding: 2")
        print("  Content: 0")
        print("  Placeholders: 2")
        print("  Decoration: 1")
        print("")
        print("Branding Elements (will be preserved):")
        print("  - image: CompanyLogo at (0.5\", 0.3\")")
        print("    [Company logo in top-left corner]")
        print("  - text: Footer at (0.5\", 7.0\")")
        print("    Text: \"© 2024 ACME Corp - Confidential\"")
        print("")
        print("Content Placeholders (available for content):")
        print("  - TITLE: Title 1 (12.3\" x 1.5\")")
        print("  - SUBTITLE: Subtitle 2 (12.3\" x 1.0\")")
        print("")
        print("Safe Content Area:")
        print("  Position: (0.5\", 1.8\")")
        print("  Size: 12.3\" x 4.9\"")
        print("  [Content will be positioned here, avoiding branding]")
        print("")


def example_2_basic_flexible_slide():
    """
    Example 2: Create a slide from template with flexible content.

    Demonstrates: Starting from template + custom content positioning.
    """
    print("="*70)
    print("Example 2: Flexible Slide from Template")
    print("="*70)

    template_path = Path("corporate_template.pptx")

    if template_path.exists():
        builder = FlexibleSlideBuilder(template_path)

        # Create slide from template layout
        slide = builder.create_slide_from_template(layout_index=1)

        # Add content - builder automatically positions within safe areas
        builder.add_content_to_slide(slide, {
            'title': 'Q4 Results Overview',
            'content_type': 'bullets',
            'content': [
                'Revenue up 25% year-over-year',
                'Expanded to 12 new markets',
                'Launched 5 major product updates',
                'Customer satisfaction at 94%',
                'Strong pipeline for Q1 2024'
            ]
        })

        builder.save(Path("flexible_basic.pptx"))

        print("\n✓ Created: flexible_basic.pptx")
        print("  - Started from template layout")
        print("  - Preserved branding (logo, footer)")
        print("  - Positioned content in safe area")
        print("  - No overlap with branding elements\n")

    else:
        print(f"\nNote: {template_path} not found. This is a demonstration.")
        print("\nWhat would happen:")
        print("  1. Load template layout #1 (Content slide)")
        print("  2. Identify branding: logo (top-left), footer (bottom)")
        print("  3. Calculate safe area: avoiding branding zones")
        print("  4. Add title to placeholder")
        print("  5. Position bullets within safe area")
        print("  6. Result: Professional slide with branding + custom content\n")


def example_3_two_column_with_branding():
    """
    Example 3: Two-column layout that respects branding.

    Shows how custom layouts work within safe areas.
    """
    print("="*70)
    print("Example 3: Custom Two-Column Layout with Branding")
    print("="*70)

    template_path = Path("corporate_template.pptx")

    if template_path.exists():
        builder = FlexibleSlideBuilder(template_path)

        slide = builder.create_slide_from_template(layout_index=1)

        # Two-column content within safe area
        builder.add_content_to_slide(slide, {
            'title': 'Our Approach vs Traditional',
            'content_type': 'two_column',
            'content': {
                'ratio': 0.5,
                'left': (
                    "Traditional Approach:\n\n"
                    "• Locked into template layouts\n"
                    "• Limited flexibility\n"
                    "• Can't customize positioning\n"
                    "• Structure defined by template"
                ),
                'right': (
                    "PPTX Agent Approach:\n\n"
                    "• Preserves branding\n"
                    "• Flexible content positioning\n"
                    "• LLM-driven layout decisions\n"
                    "• Best of both worlds"
                )
            }
        })

        builder.save(Path("flexible_two_column.pptx"))

        print("\n✓ Created: flexible_two_column.pptx")
        print("  - Template provides: logo, colors, fonts, footer")
        print("  - Custom layout: 50/50 two-column split")
        print("  - Positioned within safe area")
        print("  - Branding preserved, content flexible\n")

    else:
        print(f"\nNote: {template_path} not found. This is a demonstration.")
        print("\nConcept:")
        print("  - Template has logo (top-left) and footer (bottom)")
        print("  - Safe area calculated: avoiding branding zones")
        print("  - Two-column layout created WITHIN safe area")
        print("  - Result: Branding + custom layout = flexible + professional\n")


def example_4_image_grid_safe_positioning():
    """
    Example 4: Image grid that doesn't overlap branding.

    Demonstrates smart positioning around template elements.
    """
    print("="*70)
    print("Example 4: Image Grid with Safe Positioning")
    print("="*70)

    template_path = Path("corporate_template.pptx")

    if template_path.exists():
        builder = FlexibleSlideBuilder(template_path)

        slide = builder.create_slide_from_template(layout_index=1)

        # Create image grid (or colored placeholders if no images)
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        from pptx.dml.color import RGBColor

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = "Product Gallery"

        # Get safe area
        safe_area = builder.get_content_safe_area(slide)

        # Create 2x3 grid within safe area
        content_top = safe_area['top'] + 1.2
        content_height = safe_area['height'] - 1.2

        rows, cols = 2, 3
        gutter = 0.2

        cell_width = (safe_area['width'] - (cols - 1) * gutter) / cols
        cell_height = (content_height - (rows - 1) * gutter) / rows

        colors = [
            (255, 182, 193), (173, 216, 230), (255, 218, 185),
            (221, 160, 221), (144, 238, 144), (255, 228, 196)
        ]

        for idx in range(6):
            row = idx // cols
            col = idx % cols

            cell_left = safe_area['left'] + col * (cell_width + gutter)
            cell_top = content_top + row * (cell_height + gutter)

            # Add colored rectangle as placeholder
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cell_left), Inches(cell_top),
                Inches(cell_width), Inches(cell_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*colors[idx])
            shape.text = f"Image {idx + 1}"

        builder.save(Path("flexible_grid.pptx"))

        print("\n✓ Created: flexible_grid.pptx")
        print("  - 2x3 image grid")
        print("  - Positioned entirely within safe area")
        print("  - No overlap with logo or footer")
        print("  - Template branding preserved\n")

    else:
        print(f"\nNote: {template_path} not found. This is a demonstration.")
        print("\nConcept:")
        print("  - Template has logo (top 1.2\") and footer (bottom 0.5\")")
        print("  - Safe area: middle 5.8\" of vertical space")
        print("  - 2x3 grid calculated to fit within safe area")
        print("  - Result: Professional gallery that respects branding\n")


def example_5_llm_custom_layout_with_template():
    """
    Example 5: LLM specifies custom layout, respects template branding.

    Shows ultimate flexibility: LLM decisions + template consistency.
    """
    print("="*70)
    print("Example 5: LLM Custom Layout + Template Branding")
    print("="*70)

    template_path = Path("corporate_template.pptx")

    if template_path.exists():
        builder = FlexibleSlideBuilder(template_path)

        slide = builder.create_slide_from_template(layout_index=1)

        # LLM could generate this layout specification
        # Positions are relative (0-1) within safe content area
        llm_layout = [
            {
                'type': 'text',
                'left': 0,
                'top': 0,
                'width': 0.55,
                'height': 0.8,
                'content': (
                    "Main Insights:\n\n"
                    "• AI-driven layout decisions\n"
                    "• Template provides consistency\n"
                    "• Content positioning is flexible\n"
                    "• Best of both approaches"
                ),
                'font_size': 16
            },
            {
                'type': 'text',
                'left': 0.6,
                'top': 0,
                'width': 0.4,
                'height': 0.35,
                'content': "Key Metric:\n+45% efficiency",
                'font_size': 20
            },
            {
                'type': 'text',
                'left': 0.6,
                'top': 0.4,
                'width': 0.4,
                'height': 0.4,
                'content': "Status:\nOn Track",
                'font_size': 16
            }
        ]

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = "LLM-Designed Layout"

        # Add custom content
        builder.add_content_to_slide(slide, {
            'content_type': 'custom',
            'content': llm_layout
        })

        builder.save(Path("flexible_llm_custom.pptx"))

        print("\n✓ Created: flexible_llm_custom.pptx")
        print("  - LLM specified exact layout")
        print("  - Positioned within safe area automatically")
        print("  - Template branding preserved")
        print("  - Ultimate flexibility + consistency\n")

    else:
        print(f"\nNote: {template_path} not found. This is a demonstration.")
        print("\nKey Insight:")
        print("  This is the ULTIMATE FLEXIBILITY:")
        print("  1. Template provides: branding, colors, fonts")
        print("  2. LLM decides: exact content positioning")
        print("  3. Builder ensures: no branding overlap")
        print("  4. Result: Intelligent + consistent presentations\n")


def example_6_working_without_template():
    """
    Example 6: What if there's no template?

    Shows graceful fallback to fully dynamic layout.
    """
    print("="*70)
    print("Example 6: Flexible Builder Without Template (Fallback)")
    print("="*70)

    print("\nWhen no template is provided:")
    print("  - FlexibleSlideBuilder falls back to blank presentation")
    print("  - Uses DynamicLayoutEngine for all positioning")
    print("  - No branding constraints")
    print("  - Full creative freedom")
    print("")
    print("This ensures PPTX Agent works in all scenarios:")
    print("  ✓ With template: Branding + flexibility")
    print("  ✓ Without template: Full dynamic layout")
    print("  ✓ Best approach chosen automatically\n")


def main():
    """Run all demonstrations."""

    print("\n" + "="*70)
    print(" Flexible Template Usage - Preserve Branding, Customize Content")
    print("="*70 + "\n")

    print("This demonstration shows PPTX Agent's template intelligence:\n")
    print("Core Principles:")
    print("  1. PARSE template slides to understand structure")
    print("  2. IDENTIFY branding vs content elements")
    print("  3. PRESERVE branding (logos, watermarks, headers/footers)")
    print("  4. POSITION content flexibly within safe areas")
    print("  5. COMBINE template consistency + content flexibility\n")

    # Run examples
    example_1_analyze_template()
    example_2_basic_flexible_slide()
    example_3_two_column_with_branding()
    example_4_image_grid_safe_positioning()
    example_5_llm_custom_layout_with_template()
    example_6_working_without_template()

    print("="*70)
    print("Summary: The Complete Flexibility Solution")
    print("="*70)
    print("""
PPTX Agent solves the template rigidity problem:

❌ Traditional Approach:
  - Locked into template layout structure
  - Can't customize positioning
  - Either use template OR create from scratch
  - No middle ground

✅ PPTX Agent Approach:
  - Parse templates to understand branding
  - Preserve brand elements (logos, colors, fonts)
  - Flexible content positioning within safe areas
  - LLM can make layout decisions
  - Best of both worlds

Technical Implementation:
-------------------------
1. TemplateSlideParser - Analyzes template elements
   • Identifies branding vs content
   • Calculates safe content areas
   • Provides element metadata

2. FlexibleSlideBuilder - Builds slides intelligently
   • Starts from template layout
   • Preserves branding elements
   • Positions content dynamically
   • Respects safe area boundaries

3. DynamicLayoutEngine - Powers flexible positioning
   • Algorithmic layout calculation
   • Content-aware suggestions
   • Custom LLM-specified layouts

Result:
-------
Professional presentations that combine:
✓ Brand consistency (from templates)
✓ Content flexibility (from AI)
✓ Intelligent positioning (avoiding branding)
✓ LLM-driven decisions (optimal layouts)

This is the innovation that sets PPTX Agent apart.
    """)


if __name__ == "__main__":
    main()
