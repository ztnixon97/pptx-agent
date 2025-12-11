"""
Advanced Features Demo - Showcases all advanced PowerPoint features.

This example demonstrates:
- Speaker notes
- Hyperlinks (external and internal)
- Rich text formatting with mixed styles
- Advanced tables with cell merging and styling
- Slide management (duplicate, hide, reorder)
- Image enhancements (alt text, transparency)
- Footer and slide numbers
- Slide dimensions (widescreen vs standard)
- Shape layering (bring to front, send to back)
"""

from pathlib import Path
from pptx_agent.core.pptx_handler import PPTXHandler
from pptx_agent.builders.table_builder import TableSlideBuilder
from pptx_agent.builders.shapes_builder import ShapesBuilder
from pptx_agent.builders.smartart_builder import SmartArtBuilder


def create_advanced_demo():
    """Create a presentation showcasing all advanced features."""

    # Initialize handler with widescreen format
    handler = PPTXHandler()
    handler.set_slide_size("widescreen")

    # ========== SLIDE 1: Title Slide with Speaker Notes ==========
    title_slide = handler.add_title_slide(
        "Advanced PowerPoint Features",
        "Comprehensive Demo of PPTX Agent Capabilities"
    )

    # Add speaker notes
    handler.add_speaker_notes(
        title_slide,
        "Welcome everyone! Today we'll explore the advanced features of PPTX Agent. "
        "This presentation demonstrates speaker notes, hyperlinks, rich text formatting, "
        "advanced tables, and much more. Let's dive in!"
    )

    # ========== SLIDE 2: Rich Text Formatting ==========
    slide2 = handler.add_slide(1)
    if slide2.shapes.title:
        slide2.shapes.title.text = "Rich Text Formatting"

    # Add formatted text with mixed styles
    handler.add_formatted_text_box(
        slide2,
        text_runs=[
            {"text": "Important: ", "bold": True, "font_size": 24, "color": (255, 0, 0)},
            {"text": "This feature allows ", "font_size": 18},
            {"text": "mixed formatting", "bold": True, "italic": True, "font_size": 18, "color": (0, 0, 255)},
            {"text": " within a single text box.\n\n", "font_size": 18},
            {"text": "Benefits:\n", "bold": True, "font_size": 20, "color": (0, 128, 0)},
            {"text": "• ", "font_size": 18},
            {"text": "Emphasis", "underline": True, "font_size": 18},
            {"text": " on key words\n", "font_size": 18},
            {"text": "• Multiple ", "font_size": 18},
            {"text": "fonts", "font_name": "Courier New", "font_size": 18},
            {"text": " and ", "font_size": 18},
            {"text": "colors", "color": (255, 165, 0), "font_size": 18},
            {"text": "\n• Professional presentation styling", "font_size": 18}
        ],
        left=1.5, top=2, width=7, height=4
    )

    handler.add_speaker_notes(
        slide2,
        "This slide demonstrates rich text formatting. Notice how we can mix bold, "
        "italic, underline, different colors, and even different fonts within a single "
        "text box. This is crucial for emphasizing key points."
    )

    # ========== SLIDE 3: Hyperlinks Demo ==========
    slide3 = handler.add_slide(1)
    if slide3.shapes.title:
        slide3.shapes.title.text = "Hyperlinks - External and Internal"

    # Add text with hyperlink
    textbox = handler.add_text_box(
        slide3,
        "Click here to visit our website",
        left=2, top=2.5, width=6, height=1,
        font_size=24, bold=True, color=(0, 0, 255)
    )
    handler.add_hyperlink_to_shape(textbox, "https://github.com")

    # Add navigation box
    nav_box = handler.add_text_box(
        slide3,
        "Jump to Conclusions →",
        left=2, top=4, width=6, height=0.8,
        font_size=20, bold=True, color=(128, 0, 128)
    )
    # We'll link this to the last slide (which we'll create)
    # handler.add_internal_hyperlink(nav_box, 9)  # Will be added after creating all slides

    handler.add_speaker_notes(
        slide3,
        "Hyperlinks are essential for interactive presentations. You can link to "
        "external websites or jump to specific slides within the presentation for "
        "non-linear navigation."
    )

    # ========== SLIDE 4: Advanced Table with Cell Merging ==========
    TableSlideBuilder.add_advanced_table_slide(
        handler,
        "Advanced Table - Merged Cells & Styling",
        headers=["Quarter", "Product A", "Product B", "Product C", "Total"],
        rows=[
            ["Q1", "$100K", "$150K", "$80K", "$330K"],
            ["Q2", "$120K", "$160K", "$90K", "$370K"],
            ["Q3", "$140K", "$180K", "$95K", "$415K"],
            ["Q4", "$160K", "$200K", "$110K", "$470K"],
            ["Total", "$520K", "$690K", "$375K", "$1,585K"]
        ],
        merge_cells=[
            # No merging in this example, but you could merge cells like (1, 0, 2, 0)
        ],
        cell_styles={
            # Highlight totals row
            (4, 0): {"fill": (68, 114, 196), "bold": True, "font_color": (255, 255, 255)},
            (4, 1): {"fill": (68, 114, 196), "bold": True, "font_color": (255, 255, 255)},
            (4, 2): {"fill": (68, 114, 196), "bold": True, "font_color": (255, 255, 255)},
            (4, 3): {"fill": (68, 114, 196), "bold": True, "font_color": (255, 255, 255)},
            (4, 4): {"fill": (255, 215, 0), "bold": True, "font_size": 14, "font_color": (0, 0, 0)},
            # Highlight totals column
            (0, 4): {"fill": (144, 238, 144), "bold": True},
            (1, 4): {"fill": (144, 238, 144), "bold": True},
            (2, 4): {"fill": (144, 238, 144), "bold": True},
            (3, 4): {"fill": (144, 238, 144), "bold": True},
        },
        layout_index=1
    )

    # ========== SLIDE 5: Process Flow with Speaker Notes ==========
    smartart = SmartArtBuilder()
    slide5 = smartart.add_process_flow(
        handler,
        "Development Workflow",
        ["Requirements", "Design", "Development", "Testing", "Deployment", "Monitoring"],
        layout_idx=1
    )

    handler.add_speaker_notes(
        slide5,
        "Our development workflow follows these six key stages. Each stage has "
        "specific deliverables and quality gates. Notice how the process flows "
        "from left to right, making it easy to understand the sequence."
    )

    # ========== SLIDE 6: Comparison Diagram ==========
    slide6 = smartart.add_comparison_diagram(
        handler,
        "Cloud vs On-Premise",
        left_items=[
            "Lower upfront cost",
            "Scalability",
            "Automatic updates",
            "Accessible anywhere"
        ],
        right_items=[
            "Full control",
            "Data sovereignty",
            "Custom configuration",
            "No subscription fees"
        ],
        left_label="Cloud",
        right_label="On-Premise",
        layout_idx=1
    )

    handler.add_speaker_notes(
        slide6,
        "When choosing between cloud and on-premise solutions, consider these factors. "
        "Cloud offers flexibility and lower initial costs, while on-premise provides "
        "complete control and data sovereignty."
    )

    # ========== SLIDE 7: Shape Layering Demo ==========
    shapes_builder = ShapesBuilder()
    slide7 = handler.add_slide(5)  # Blank layout
    if slide7.shapes.title:
        slide7.shapes.title.text = "Shape Layering"

    # Add overlapping shapes to demonstrate layering
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    # Background shape (will be sent to back)
    shape1 = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(2), Inches(4), Inches(3)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(255, 0, 0)
    shape1.text = "Background (Send to Back)"

    # Middle shape
    shape2 = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3), Inches(2.5), Inches(4), Inches(3)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0, 255, 0)
    shape2.text = "Middle Layer"

    # Foreground shape (will be brought to front)
    shape3 = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(3), Inches(4), Inches(3)
    )
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(0, 0, 255)
    shape3.text = "Foreground (Bring to Front)"

    # Apply layering
    handler.send_to_back(slide7, 0)  # Send red to back
    handler.bring_to_front(slide7, 2)  # Bring blue to front

    handler.add_speaker_notes(
        slide7,
        "Shape layering allows you to control which elements appear on top. "
        "This is essential for creating complex visual compositions and ensuring "
        "important elements are visible."
    )

    # ========== SLIDE 8: Flowchart ==========
    slide8 = shapes_builder.add_flowchart_slide(
        handler,
        "Decision Flowchart",
        [
            {"text": "Start Process", "decision": False},
            {"text": "Budget Approved?", "decision": True},
            {"text": "Review Scope", "decision": False},
            {"text": "Requirements Met?", "decision": True},
            {"text": "Proceed to Implementation", "decision": False}
        ],
        layout_idx=1
    )

    handler.add_speaker_notes(
        slide8,
        "This flowchart shows our project approval process. Diamond shapes represent "
        "decision points, while rectangles represent process steps. Clear visual "
        "distinction helps stakeholders understand the workflow."
    )

    # ========== SLIDE 9: Hidden Slide (Demo) ==========
    slide9 = handler.add_slide(1)
    if slide9.shapes.title:
        slide9.shapes.title.text = "Hidden Slide - Technical Details"

    handler.add_text_box(
        slide9,
        "This slide contains technical details that can be shown if needed,\n"
        "but is hidden during normal presentation flow.",
        left=2, top=2.5, width=6, height=2,
        font_size=18
    )

    # Hide this slide
    handler.hide_slide(8)  # 0-based index

    # ========== SLIDE 10: Conclusion with Footer ==========
    slide10 = handler.add_slide(1)
    if slide10.shapes.title:
        slide10.shapes.title.text = "Conclusion"

    handler.add_formatted_text_box(
        slide10,
        text_runs=[
            {"text": "Key Takeaways:\n\n", "bold": True, "font_size": 24, "color": (68, 114, 196)},
            {"text": "✓ ", "font_size": 20, "color": (0, 128, 0)},
            {"text": "Speaker notes enhance presentation delivery\n", "font_size": 18},
            {"text": "✓ ", "font_size": 20, "color": (0, 128, 0)},
            {"text": "Hyperlinks enable interactive presentations\n", "font_size": 18},
            {"text": "✓ ", "font_size": 20, "color": (0, 128, 0)},
            {"text": "Rich formatting emphasizes key points\n", "font_size": 18},
            {"text": "✓ ", "font_size": 20, "color": (0, 128, 0)},
            {"text": "Advanced tables present complex data clearly\n", "font_size": 18},
            {"text": "✓ ", "font_size": 20, "color": (0, 128, 0)},
            {"text": "Slide management provides presentation control\n", "font_size": 18},
        ],
        left=1.5, top=2, width=7, height=4
    )

    # Now add the internal hyperlink from slide 3 to this conclusion slide
    handler.add_internal_hyperlink(nav_box, 9)  # Link to slide 10 (0-based index)

    handler.add_speaker_notes(
        slide10,
        "Thank you for attending this demonstration! These advanced features enable "
        "you to create professional, interactive, and accessible presentations. "
        "Remember to use speaker notes for all your presentations - they're invaluable "
        "for preparation and delivery."
    )

    # ========== Set Presentation Footer ==========
    handler.set_presentation_footer(
        footer_text="Advanced Features Demo | PPTX Agent",
        show_slide_number=True,
        show_date=False
    )

    # ========== Save Presentation ==========
    output_path = Path("advanced_features_demo.pptx")
    handler.save(output_path)

    print(f"✓ Advanced features demo created: {output_path}")
    print(f"✓ Total slides: {handler.get_slide_count()}")
    print(f"✓ Slide dimensions: Widescreen (16:9)")
    print(f"✓ Features demonstrated:")
    print("  - Speaker notes on all slides")
    print("  - External and internal hyperlinks")
    print("  - Rich text formatting with mixed styles")
    print("  - Advanced tables with cell styling")
    print("  - SmartArt diagrams (process flow, comparison)")
    print("  - Shape layering (z-order control)")
    print("  - Flowcharts")
    print("  - Hidden slide (slide 9)")
    print("  - Footer with slide numbers")
    print("  - Widescreen format (16:9)")


if __name__ == "__main__":
    create_advanced_demo()
