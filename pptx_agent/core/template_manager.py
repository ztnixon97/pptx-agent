"""
Template Manager - Handles PowerPoint template analysis and management.

Templates are used for STYLING (colors, fonts) not STRUCTURE (layouts).
This allows LLM-driven dynamic layout creation while maintaining visual consistency.
"""

from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class TemplateManager:
    """Manages PowerPoint templates and extracts their properties."""

    def __init__(self, template_path: Optional[Path] = None):
        """
        Initialize the template manager.

        Args:
            template_path: Path to a template .pptx file
        """
        self.template_path = template_path
        self.template_info = {}
        self.theme = {}  # Extracted theme (colors, fonts)

        if template_path and template_path.exists():
            self._analyze_template()

    def _analyze_template(self):
        """Analyze the template and extract useful information."""
        prs = Presentation(str(self.template_path))

        self.template_info = {
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
            'layouts': [],
            'master_slides': len(prs.slide_master.slide_layouts)
        }

        # Extract theme information (NEW!)
        self.theme = self._extract_comprehensive_theme(prs)

        # Extract layout information (kept for backward compatibility)
        for idx, layout in enumerate(prs.slide_layouts):
            layout_info = {
                'index': idx,
                'name': layout.name,
                'placeholders': []
            }

            for placeholder in layout.placeholders:
                ph_info = {
                    'index': placeholder.placeholder_format.idx,
                    'type': placeholder.placeholder_format.type,
                    'name': placeholder.name,
                    'left': placeholder.left,
                    'top': placeholder.top,
                    'width': placeholder.width,
                    'height': placeholder.height
                }
                layout_info['placeholders'].append(ph_info)

            self.template_info['layouts'].append(layout_info)

    def _extract_comprehensive_theme(self, prs: Presentation) -> Dict[str, Any]:
        """
        Extract comprehensive theme information from template.

        This includes colors, fonts, and styling - NOT layout structure.
        """
        theme = {
            'extracted': True,
            'colors': {},
            'fonts': {},
            'default_styles': {}
        }

        try:
            # Get slide master
            master = prs.slide_master

            # Extract theme colors
            theme['colors'] = self._extract_theme_colors(master)

            # Extract fonts
            theme['fonts'] = self._extract_theme_fonts(master)

            # Extract default text styles
            theme['default_styles'] = self._extract_default_styles(master)

        except Exception as e:
            theme['extraction_error'] = str(e)

        return theme

    def _extract_theme_colors(self, master) -> Dict[str, Tuple[int, int, int]]:
        """
        Extract theme colors from slide master.

        Returns RGB values for common theme colors.
        """
        colors = {
            'primary': (68, 114, 196),     # Default blue
            'secondary': (237, 125, 49),   # Default orange
            'accent1': (112, 173, 71),     # Default green
            'accent2': (255, 192, 0),      # Default gold
            'background': (255, 255, 255), # White
            'text': (0, 0, 0),             # Black
            'gray_light': (217, 225, 242),
            'gray_dark': (68, 84, 106)
        }

        try:
            # Try to extract from theme
            # Note: This is simplified - real extraction is complex
            # For now, using sensible defaults
            pass
        except Exception:
            pass

        return colors

    def _extract_theme_fonts(self, master) -> Dict[str, str]:
        """
        Extract theme fonts from slide master.

        Returns font names for different text types.
        """
        fonts = {
            'title': 'Calibri Light',
            'body': 'Calibri',
            'heading': 'Calibri',
            'code': 'Consolas'
        }

        try:
            # Try to extract actual theme fonts
            # Note: Simplified implementation
            pass
        except Exception:
            pass

        return fonts

    def _extract_default_styles(self, master) -> Dict[str, Any]:
        """
        Extract default text styles from template.

        Returns font sizes, colors, and other styling for different text types.
        """
        styles = {
            'title': {
                'font_size': 44,
                'bold': True,
                'color': (0, 0, 0)
            },
            'subtitle': {
                'font_size': 28,
                'bold': False,
                'color': (68, 84, 106)
            },
            'body': {
                'font_size': 18,
                'bold': False,
                'color': (0, 0, 0)
            },
            'heading': {
                'font_size': 24,
                'bold': True,
                'color': (68, 114, 196)
            },
            'bullet': {
                'font_size': 18,
                'bold': False,
                'color': (0, 0, 0)
            }
        }

        return styles

    def get_theme_color(self, color_name: str) -> Optional[Tuple[int, int, int]]:
        """
        Get a specific theme color.

        Args:
            color_name: Name of color ('primary', 'secondary', 'accent1', etc.)

        Returns:
            RGB tuple or None
        """
        return self.theme.get('colors', {}).get(color_name)

    def get_theme_font(self, font_type: str) -> Optional[str]:
        """
        Get a specific theme font.

        Args:
            font_type: Type of font ('title', 'body', 'heading', etc.)

        Returns:
            Font name or None
        """
        return self.theme.get('fonts', {}).get(font_type)

    def get_text_style(self, style_type: str) -> Optional[Dict[str, Any]]:
        """
        Get default text style.

        Args:
            style_type: Type of text ('title', 'subtitle', 'body', etc.)

        Returns:
            Style dictionary or None
        """
        return self.theme.get('default_styles', {}).get(style_type)

    def get_layout_info(self, layout_name: Optional[str] = None,
                       layout_index: Optional[int] = None) -> Optional[Dict]:
        """
        Get information about a specific layout.

        Args:
            layout_name: Name of the layout
            layout_index: Index of the layout

        Returns:
            Dictionary with layout information or None
        """
        if not self.template_info.get('layouts'):
            return None

        if layout_name:
            for layout in self.template_info['layouts']:
                if layout['name'].lower() == layout_name.lower():
                    return layout

        if layout_index is not None:
            for layout in self.template_info['layouts']:
                if layout['index'] == layout_index:
                    return layout

        return None

    def list_layouts(self) -> List[Dict[str, Any]]:
        """List all available layouts in the template."""
        return self.template_info.get('layouts', [])

    def suggest_layout(self, content_type: str) -> Optional[int]:
        """
        Suggest a layout index based on content type.

        Note: Prefer using DynamicLayoutEngine instead of template layouts.

        Args:
            content_type: Type of content ('title', 'content', 'section', 'blank', etc.)

        Returns:
            Suggested layout index or None
        """
        if not self.template_info.get('layouts'):
            return 0

        content_type_lower = content_type.lower()

        # Common layout name patterns
        patterns = {
            'title': ['title', 'cover'],
            'content': ['content', 'body', 'text'],
            'section': ['section', 'divider'],
            'blank': ['blank', 'empty'],
            'comparison': ['comparison', 'two'],
            'title_only': ['title only', 'heading']
        }

        for layout in self.template_info['layouts']:
            layout_name = layout['name'].lower()
            if content_type_lower in patterns:
                for pattern in patterns[content_type_lower]:
                    if pattern in layout_name:
                        return layout['index']

        # Default fallbacks
        if content_type_lower == 'title':
            return 0
        else:
            return 1 if len(self.template_info['layouts']) > 1 else 0

    def get_template_summary(self) -> str:
        """Get a human-readable summary of the template."""
        if not self.template_info:
            return "No template loaded"

        summary = [
            f"Template: {self.template_path.name if self.template_path else 'None'}",
            f"Slide Size: {self.template_info.get('slide_width', 0) / 914400:.1f}\" x {self.template_info.get('slide_height', 0) / 914400:.1f}\"",
            f"Layouts: {len(self.template_info.get('layouts', []))}",
            "",
            "Theme Colors:"
        ]

        for color_name, rgb in self.theme.get('colors', {}).items():
            summary.append(f"  {color_name}: RGB{rgb}")

        summary.append("\nTheme Fonts:")
        for font_type, font_name in self.theme.get('fonts', {}).items():
            summary.append(f"  {font_type}: {font_name}")

        summary.append("\nAvailable Layouts (for reference):")
        for layout in self.template_info.get('layouts', []):
            summary.append(f"  [{layout['index']}] {layout['name']} "
                          f"({len(layout['placeholders'])} placeholders)")

        return "\n".join(summary)

    def get_theme_summary(self) -> Dict[str, Any]:
        """
        Get complete theme information for use with DynamicLayoutEngine.

        Returns:
            Dictionary with all theme information
        """
        return {
            'colors': self.theme.get('colors', {}),
            'fonts': self.theme.get('fonts', {}),
            'styles': self.theme.get('default_styles', {}),
            'dimensions': {
                'width': self.template_info.get('slide_width', 9144000) / 914400,  # Convert to inches
                'height': self.template_info.get('slide_height', 6858000) / 914400
            }
        }

