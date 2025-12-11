"""
Template Slide Parser - Analyze template slides to identify and classify elements.

This module enables intelligent template usage by:
- Identifying all elements on template slides (shapes, text, images, etc.)
- Classifying elements as 'branding' (preserve) vs 'content' (modifiable)
- Extracting element properties and positions
- Enabling preservation of branding while allowing content flexibility
"""

from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE


@dataclass
class SlideElement:
    """Represents an element found on a template slide."""
    element_id: int
    element_type: str  # 'text', 'image', 'shape', 'table', 'chart', 'group'
    classification: str  # 'branding', 'content', 'placeholder', 'decoration'
    left: float  # inches
    top: float  # inches
    width: float  # inches
    height: float  # inches
    is_placeholder: bool
    placeholder_type: Optional[str]
    text_content: str
    name: str
    has_image: bool
    shape_type: Optional[str]
    z_index: int

    @property
    def bounds(self) -> Dict[str, float]:
        """Get element bounds."""
        return {
            'left': self.left,
            'top': self.top,
            'right': self.left + self.width,
            'bottom': self.top + self.height,
            'width': self.width,
            'height': self.height
        }


class TemplateSlideParser:
    """Parse template slides to identify and classify elements."""

    def __init__(self, template_path: Path):
        """
        Initialize the template slide parser.

        Args:
            template_path: Path to the template .pptx file
        """
        self.template_path = template_path
        self.presentation = Presentation(str(template_path))
        self.slide_layouts = {}  # layout_index -> parsed layout info

    def parse_layout(self, layout_index: int) -> Dict[str, Any]:
        """
        Parse a specific slide layout to understand its elements.

        Args:
            layout_index: Index of the layout to parse

        Returns:
            Dictionary with layout information and elements
        """
        if layout_index >= len(self.presentation.slide_layouts):
            raise ValueError(f"Layout index {layout_index} out of range")

        layout = self.presentation.slide_layouts[layout_index]

        # Parse all elements on the layout
        elements = []
        for idx, shape in enumerate(layout.shapes):
            element = self._parse_shape(shape, idx)
            if element:
                elements.append(element)

        # Classify elements
        branding_elements = [e for e in elements if e.classification == 'branding']
        content_elements = [e for e in elements if e.classification == 'content']
        placeholder_elements = [e for e in elements if e.classification == 'placeholder']
        decoration_elements = [e for e in elements if e.classification == 'decoration']

        layout_info = {
            'index': layout_index,
            'name': layout.name,
            'all_elements': elements,
            'branding_elements': branding_elements,
            'content_elements': content_elements,
            'placeholder_elements': placeholder_elements,
            'decoration_elements': decoration_elements,
            'slide_width': self.presentation.slide_width / 914400,  # Convert to inches
            'slide_height': self.presentation.slide_height / 914400
        }

        # Cache the parsed layout
        self.slide_layouts[layout_index] = layout_info

        return layout_info

    def _parse_shape(self, shape, index: int) -> Optional[SlideElement]:
        """
        Parse a single shape and extract its properties.

        Args:
            shape: The shape to parse
            index: Shape index for z-order

        Returns:
            SlideElement or None if shape should be ignored
        """
        try:
            # Extract basic properties
            left = shape.left / 914400  # Convert EMU to inches
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400

            # Determine element type
            element_type = self._get_element_type(shape)

            # Extract text content
            text_content = ""
            if hasattr(shape, 'text_frame'):
                text_content = shape.text_frame.text
            elif hasattr(shape, 'text'):
                text_content = shape.text

            # Check if placeholder
            is_placeholder = False
            placeholder_type = None
            if shape.is_placeholder:
                is_placeholder = True
                try:
                    placeholder_type = str(shape.placeholder_format.type)
                except:
                    placeholder_type = "unknown"

            # Check if contains image
            has_image = element_type == 'image'

            # Get shape type
            shape_type = None
            if hasattr(shape, 'shape_type'):
                try:
                    shape_type = str(shape.shape_type)
                except:
                    pass

            # Classify the element
            classification = self._classify_element(
                element_type, is_placeholder, placeholder_type,
                text_content, left, top, width, height, shape
            )

            return SlideElement(
                element_id=index,
                element_type=element_type,
                classification=classification,
                left=left,
                top=top,
                width=width,
                height=height,
                is_placeholder=is_placeholder,
                placeholder_type=placeholder_type,
                text_content=text_content,
                name=shape.name if hasattr(shape, 'name') else f"Shape_{index}",
                has_image=has_image,
                shape_type=shape_type,
                z_index=index
            )

        except Exception as e:
            # Skip shapes that can't be parsed
            return None

    def _get_element_type(self, shape) -> str:
        """Determine the element type from the shape."""
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type

            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                return 'image'
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                return 'table'
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                return 'chart'
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                return 'group'
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, 'text_frame'):
                return 'text'
            elif shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
                return 'shape'

        # Fallback: check for text frame
        if hasattr(shape, 'text_frame'):
            return 'text'

        return 'shape'

    def _classify_element(self, element_type: str, is_placeholder: bool,
                         placeholder_type: Optional[str], text_content: str,
                         left: float, top: float, width: float, height: float,
                         shape) -> str:
        """
        Classify element as branding, content, placeholder, or decoration.

        Classification rules:
        - Branding: Logos, company names, watermarks, persistent headers/footers
        - Content: Placeholders, main content areas
        - Placeholder: Empty or template text placeholders
        - Decoration: Decorative shapes, lines, backgrounds
        """
        # Placeholders are usually content areas
        if is_placeholder:
            # Title and content placeholders
            if placeholder_type and ('TITLE' in str(placeholder_type).upper() or
                                     'BODY' in str(placeholder_type).upper() or
                                     'OBJECT' in str(placeholder_type).upper()):
                return 'placeholder'
            # Footer, slide number, date placeholders are branding
            elif placeholder_type and ('FOOTER' in str(placeholder_type).upper() or
                                       'NUMBER' in str(placeholder_type).upper() or
                                       'DATE' in str(placeholder_type).upper()):
                return 'branding'
            else:
                return 'placeholder'

        # Images in specific locations are likely logos (top corners, bottom corners)
        if element_type == 'image':
            slide_width = self.presentation.slide_width / 914400
            slide_height = self.presentation.slide_height / 914400

            # Top-left corner (likely logo)
            if left < 1.5 and top < 1.5:
                return 'branding'

            # Top-right corner
            if left > (slide_width - width - 1.5) and top < 1.5:
                return 'branding'

            # Bottom corners (watermarks)
            if top > (slide_height - height - 1.0):
                return 'branding'

            # Small images are likely logos
            if width < 2.0 and height < 2.0:
                return 'branding'

            # Otherwise, content image
            return 'content'

        # Text analysis for branding
        if element_type == 'text' and text_content:
            text_lower = text_content.lower()

            # Company name indicators
            branding_keywords = ['company', 'inc', 'llc', 'corp', 'ltd',
                               'copyright', '©', '®', '™', 'confidential',
                               'proprietary', 'all rights reserved']

            if any(keyword in text_lower for keyword in branding_keywords):
                return 'branding'

            # Small text at edges (headers/footers)
            slide_height = self.presentation.slide_height / 914400
            if (top < 0.75 or top > slide_height - 1.0) and len(text_content) < 100:
                return 'branding'

        # Small decorative shapes
        if element_type == 'shape':
            # Very small shapes are likely decorative
            if width < 0.5 or height < 0.5:
                return 'decoration'

            # Thin horizontal/vertical lines
            if (width > 5.0 and height < 0.2) or (height > 3.0 and width < 0.2):
                return 'decoration'

            # No text and small size = decoration
            if not text_content and width < 2.0 and height < 2.0:
                return 'decoration'

        # Groups are often branding (logo + text combinations)
        if element_type == 'group':
            slide_width = self.presentation.slide_width / 914400
            # Groups in corners are likely branding
            if left < 2.0 or left > (slide_width - width - 2.0):
                return 'branding'

        # Default: content
        return 'content'

    def get_content_safe_area(self, layout_index: int) -> Dict[str, float]:
        """
        Calculate the safe area for content (avoiding branding elements).

        Args:
            layout_index: Index of the layout

        Returns:
            Dictionary with left, top, width, height for safe content area
        """
        layout_info = self.parse_layout(layout_index)

        slide_width = layout_info['slide_width']
        slide_height = layout_info['slide_height']

        # Start with full slide
        safe_left = 0.5  # Default margin
        safe_top = 0.5
        safe_right = slide_width - 0.5
        safe_bottom = slide_height - 0.5

        # Adjust for branding elements
        for element in layout_info['branding_elements']:
            bounds = element.bounds

            # If branding is at top, move safe area down
            if bounds['top'] < 1.5 and bounds['bottom'] < slide_height / 3:
                safe_top = max(safe_top, bounds['bottom'] + 0.2)

            # If branding is at bottom, move safe area up
            if bounds['bottom'] > slide_height - 1.0:
                safe_bottom = min(safe_bottom, bounds['top'] - 0.2)

            # If branding is at left, move safe area right
            if bounds['left'] < 1.5 and bounds['right'] < slide_width / 4:
                safe_left = max(safe_left, bounds['right'] + 0.2)

            # If branding is at right, move safe area left
            if bounds['right'] > slide_width - 1.5:
                safe_right = min(safe_right, bounds['left'] - 0.2)

        return {
            'left': safe_left,
            'top': safe_top,
            'width': safe_right - safe_left,
            'height': safe_bottom - safe_top,
            'right': safe_right,
            'bottom': safe_bottom
        }

    def get_layout_summary(self, layout_index: int) -> str:
        """
        Get a human-readable summary of a layout's elements.

        Args:
            layout_index: Index of the layout

        Returns:
            Formatted string describing the layout
        """
        layout_info = self.parse_layout(layout_index)

        lines = [
            f"Layout: {layout_info['name']} (Index {layout_index})",
            f"Dimensions: {layout_info['slide_width']:.1f}\" x {layout_info['slide_height']:.1f}\"",
            "",
            f"Total Elements: {len(layout_info['all_elements'])}",
            f"  Branding: {len(layout_info['branding_elements'])}",
            f"  Content: {len(layout_info['content_elements'])}",
            f"  Placeholders: {len(layout_info['placeholder_elements'])}",
            f"  Decoration: {len(layout_info['decoration_elements'])}",
            ""
        ]

        if layout_info['branding_elements']:
            lines.append("Branding Elements (will be preserved):")
            for elem in layout_info['branding_elements']:
                lines.append(f"  - {elem.element_type}: {elem.name} at ({elem.left:.1f}\", {elem.top:.1f}\")")
                if elem.text_content:
                    preview = elem.text_content[:50] + "..." if len(elem.text_content) > 50 else elem.text_content
                    lines.append(f"    Text: \"{preview}\"")
            lines.append("")

        if layout_info['placeholder_elements']:
            lines.append("Content Placeholders (available for content):")
            for elem in layout_info['placeholder_elements']:
                lines.append(f"  - {elem.placeholder_type}: {elem.name} ({elem.width:.1f}\" x {elem.height:.1f}\")")
            lines.append("")

        safe_area = self.get_content_safe_area(layout_index)
        lines.append(f"Safe Content Area:")
        lines.append(f"  Position: ({safe_area['left']:.1f}\", {safe_area['top']:.1f}\")")
        lines.append(f"  Size: {safe_area['width']:.1f}\" x {safe_area['height']:.1f}\"")

        return "\n".join(lines)

    def list_all_layouts(self) -> List[Dict[str, Any]]:
        """
        Parse and return information about all layouts in the template.

        Returns:
            List of layout info dictionaries
        """
        layouts = []
        for idx in range(len(self.presentation.slide_layouts)):
            layout_info = self.parse_layout(idx)
            layouts.append(layout_info)
        return layouts
