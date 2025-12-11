"""
Flexible Slide Builder - Build slides from templates while preserving branding.

This module enables intelligent template-based slide creation:
- Starts from an actual template slide (not blank)
- Preserves branding elements (logos, watermarks, headers/footers)
- Allows flexible modification of content areas
- Combines template styling with dynamic content positioning
"""

from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import copy

from .template_slide_parser import TemplateSlideParser, SlideElement
from .pptx_handler import PPTXHandler
from .dynamic_layout_engine import DynamicLayoutEngine, ElementBox


class FlexibleSlideBuilder:
    """
    Build slides from templates while preserving branding and allowing content flexibility.

    This builder:
    1. Analyzes template layouts to identify branding vs content
    2. Preserves branding elements in their original positions
    3. Allows dynamic positioning of content within safe areas
    4. Combines template consistency with content flexibility
    """

    def __init__(self, template_path: Optional[Path] = None):
        """
        Initialize the flexible slide builder.

        Args:
            template_path: Path to template .pptx file
        """
        self.template_path = template_path
        self.parser = None
        self.handler = None
        self.layout_engine = None

        if template_path and template_path.exists():
            self.parser = TemplateSlideParser(template_path)
            self.handler = PPTXHandler(template_path)

            # Initialize layout engine with template dimensions
            prs = Presentation(str(template_path))
            slide_width = prs.slide_width / 914400  # Convert to inches
            slide_height = prs.slide_height / 914400
            self.layout_engine = DynamicLayoutEngine(slide_width, slide_height)

    def create_slide_from_template(self, layout_index: int = 1,
                                   preserve_branding: bool = True,
                                   preserve_decorations: bool = True) -> Any:
        """
        Create a slide based on a template layout.

        Args:
            layout_index: Index of the template layout to use
            preserve_branding: Whether to keep branding elements
            preserve_decorations: Whether to keep decorative elements

        Returns:
            The created slide
        """
        if not self.handler:
            raise ValueError("No template loaded. Provide template_path when initializing.")

        # Parse the layout to understand elements
        layout_info = self.parser.parse_layout(layout_index)

        # Create slide using the template layout
        slide = self.handler.add_slide(layout_index)

        # Store layout info for this slide
        slide._layout_info = layout_info

        return slide

    def get_content_safe_area(self, slide) -> Dict[str, float]:
        """
        Get the safe area for content on this slide (avoiding branding).

        Args:
            slide: The slide to analyze

        Returns:
            Dictionary with left, top, width, height for safe content area
        """
        if hasattr(slide, '_layout_info'):
            layout_info = slide._layout_info

            slide_width = layout_info['slide_width']
            slide_height = layout_info['slide_height']

            safe_left = 0.5
            safe_top = 0.5
            safe_right = slide_width - 0.5
            safe_bottom = slide_height - 0.5

            # Adjust for branding elements
            for element in layout_info['branding_elements']:
                bounds = element.bounds

                if bounds['top'] < 1.5 and bounds['bottom'] < slide_height / 3:
                    safe_top = max(safe_top, bounds['bottom'] + 0.2)

                if bounds['bottom'] > slide_height - 1.0:
                    safe_bottom = min(safe_bottom, bounds['top'] - 0.2)

                if bounds['left'] < 1.5 and bounds['right'] < slide_width / 4:
                    safe_left = max(safe_left, bounds['right'] + 0.2)

                if bounds['right'] > slide_width - 1.5:
                    safe_right = min(safe_right, bounds['left'] - 0.2)

            return {
                'left': safe_left,
                'top': safe_top,
                'width': safe_right - safe_left,
                'height': safe_bottom - safe_top,
                'right': safe_right,
                'bottom': safe_bottom
            }
        else:
            # Fallback: use template parser
            layout_index = slide.slide_layout.slide_layouts.index(slide.slide_layout)
            return self.parser.get_content_safe_area(layout_index)

    def add_content_to_slide(self, slide, content_spec: Dict[str, Any]):
        """
        Add content to a slide, automatically positioning within safe areas.

        Args:
            slide: The slide to add content to
            content_spec: Dictionary describing content to add
                {
                    'title': str,
                    'content_type': str ('text', 'bullets', 'two_column', 'grid', etc.),
                    'content': Any (varies by content_type),
                    'layout_style': Optional[str] ('auto', 'custom')
                }
        """
        safe_area = self.get_content_safe_area(slide)

        content_type = content_spec.get('content_type', 'text')
        title = content_spec.get('title')
        content = content_spec.get('content')

        # Handle title
        if title and slide.shapes.title:
            slide.shapes.title.text = title

        # Position content based on type
        if content_type == 'text':
            self._add_text_content(slide, content, safe_area)
        elif content_type == 'bullets':
            self._add_bullet_content(slide, content, safe_area)
        elif content_type == 'two_column':
            self._add_two_column_content(slide, content, safe_area)
        elif content_type == 'grid':
            self._add_grid_content(slide, content, safe_area)
        elif content_type == 'image':
            self._add_image_content(slide, content, safe_area)
        elif content_type == 'custom':
            self._add_custom_content(slide, content, safe_area)

    def _add_text_content(self, slide, content: str, safe_area: Dict[str, float]):
        """Add text content within safe area."""
        # Skip title area
        content_top = safe_area['top'] + 1.2
        content_height = safe_area['height'] - 1.2

        self.handler.add_text_box(
            slide,
            content,
            left=safe_area['left'],
            top=content_top,
            width=safe_area['width'],
            height=content_height,
            font_size=18
        )

    def _add_bullet_content(self, slide, bullets: List[str], safe_area: Dict[str, float]):
        """Add bullet points within safe area."""
        content_top = safe_area['top'] + 1.2
        content_height = safe_area['height'] - 1.2

        self.handler.add_bullet_slide(
            slide=slide,
            bullets=bullets,
            left=safe_area['left'],
            top=content_top,
            width=safe_area['width'],
            height=content_height
        )

    def _add_two_column_content(self, slide, content: Dict[str, Any], safe_area: Dict[str, float]):
        """Add two-column layout within safe area."""
        content_top = safe_area['top'] + 1.2
        content_height = safe_area['height'] - 1.2

        column_ratio = content.get('ratio', 0.5)
        gutter = 0.3

        left_width = (safe_area['width'] - gutter) * column_ratio
        right_width = safe_area['width'] - left_width - gutter

        # Left column
        left_content = content.get('left', '')
        if left_content:
            self.handler.add_text_box(
                slide,
                left_content,
                left=safe_area['left'],
                top=content_top,
                width=left_width,
                height=content_height,
                font_size=16
            )

        # Right column
        right_content = content.get('right', '')
        if right_content:
            self.handler.add_text_box(
                slide,
                right_content,
                left=safe_area['left'] + left_width + gutter,
                top=content_top,
                width=right_width,
                height=content_height,
                font_size=16
            )

    def _add_grid_content(self, slide, content: Dict[str, Any], safe_area: Dict[str, float]):
        """Add grid layout within safe area."""
        content_top = safe_area['top'] + 1.2
        content_height = safe_area['height'] - 1.2

        rows = content.get('rows', 2)
        cols = content.get('cols', 2)
        items = content.get('items', [])

        gutter = 0.2
        cell_width = (safe_area['width'] - (cols - 1) * gutter) / cols
        cell_height = (content_height - (rows - 1) * gutter) / rows

        for idx, item in enumerate(items[:rows * cols]):
            row = idx // cols
            col = idx % cols

            cell_left = safe_area['left'] + col * (cell_width + gutter)
            cell_top = content_top + row * (cell_height + gutter)

            if isinstance(item, str):
                # Text item
                self.handler.add_text_box(
                    slide,
                    item,
                    left=cell_left,
                    top=cell_top,
                    width=cell_width,
                    height=cell_height,
                    font_size=14
                )
            elif isinstance(item, dict) and 'image' in item:
                # Image item
                self.handler.add_image(
                    slide,
                    Path(item['image']),
                    left=cell_left,
                    top=cell_top,
                    width=cell_width,
                    height=cell_height
                )

    def _add_image_content(self, slide, content: Dict[str, Any], safe_area: Dict[str, float]):
        """Add image within safe area."""
        image_path = content.get('path')
        if not image_path:
            return

        position = content.get('position', 'center')
        size = content.get('size', 'medium')

        content_top = safe_area['top'] + 1.2
        content_height = safe_area['height'] - 1.2

        # Calculate image size
        if size == 'small':
            img_width = safe_area['width'] * 0.4
            img_height = content_height * 0.5
        elif size == 'large':
            img_width = safe_area['width'] * 0.8
            img_height = content_height * 0.8
        else:  # medium
            img_width = safe_area['width'] * 0.6
            img_height = content_height * 0.6

        # Calculate position
        if position == 'left':
            img_left = safe_area['left']
        elif position == 'right':
            img_left = safe_area['right'] - img_width
        else:  # center
            img_left = safe_area['left'] + (safe_area['width'] - img_width) / 2

        img_top = content_top + (content_height - img_height) / 2

        self.handler.add_image(
            slide,
            Path(image_path),
            left=img_left,
            top=img_top,
            width=img_width,
            height=img_height
        )

    def _add_custom_content(self, slide, content: List[Dict[str, Any]], safe_area: Dict[str, float]):
        """
        Add custom positioned content within safe area.

        Content is a list of element specifications with relative positions.
        """
        for element_spec in content:
            elem_type = element_spec.get('type')

            # Convert relative positions (0-1) to absolute within safe area
            rel_left = element_spec.get('left', 0)
            rel_top = element_spec.get('top', 0)
            rel_width = element_spec.get('width', 1.0)
            rel_height = element_spec.get('height', 0.5)

            content_top = safe_area['top'] + 1.2
            content_height = safe_area['height'] - 1.2

            abs_left = safe_area['left'] + rel_left * safe_area['width']
            abs_top = content_top + rel_top * content_height
            abs_width = rel_width * safe_area['width']
            abs_height = rel_height * content_height

            if elem_type == 'text':
                self.handler.add_text_box(
                    slide,
                    element_spec.get('content', ''),
                    left=abs_left,
                    top=abs_top,
                    width=abs_width,
                    height=abs_height,
                    font_size=element_spec.get('font_size', 16)
                )
            elif elem_type == 'image':
                self.handler.add_image(
                    slide,
                    Path(element_spec.get('path')),
                    left=abs_left,
                    top=abs_top,
                    width=abs_width,
                    height=abs_height
                )

    def clear_placeholder_content(self, slide):
        """
        Clear content from placeholder areas while preserving branding.

        Args:
            slide: The slide to clear
        """
        if not hasattr(slide, '_layout_info'):
            return

        layout_info = slide._layout_info

        # Identify placeholder elements to clear
        for element in layout_info['placeholder_elements']:
            # Find corresponding shape on slide
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.clear()
                        elif hasattr(shape, 'text'):
                            shape.text = ""
                    except:
                        pass

    def get_layout_summary(self, layout_index: int) -> str:
        """
        Get summary of a template layout.

        Args:
            layout_index: Index of the layout

        Returns:
            Human-readable summary
        """
        if not self.parser:
            return "No template loaded"

        return self.parser.get_layout_summary(layout_index)

    def list_available_layouts(self) -> List[Dict[str, Any]]:
        """
        List all available layouts from the template.

        Returns:
            List of layout information dictionaries
        """
        if not self.parser:
            return []

        return self.parser.list_all_layouts()

    def save(self, output_path: Path):
        """
        Save the presentation.

        Args:
            output_path: Path to save the .pptx file
        """
        if self.handler:
            self.handler.save(output_path)
        else:
            raise ValueError("No handler available. Template not loaded.")
