"""
PowerPoint Handler - Core utilities for working with PowerPoint presentations.
"""

from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


class PPTXHandler:
    """Handles PowerPoint presentation creation and manipulation."""

    def __init__(self, template_path: Optional[Path] = None):
        """
        Initialize the PPTX handler.

        Args:
            template_path: Optional path to a template .pptx file
        """
        if template_path and template_path.exists():
            self.prs = Presentation(str(template_path))
            self.template_path = template_path
        else:
            self.prs = Presentation()
            self.template_path = None

        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def get_layout_by_name(self, name: str):
        """Get a slide layout by name."""
        for layout in self.prs.slide_layouts:
            if layout.name.lower() == name.lower():
                return layout
        return None

    def get_layout_by_index(self, index: int):
        """Get a slide layout by index."""
        try:
            return self.prs.slide_layouts[index]
        except IndexError:
            return self.prs.slide_layouts[0]

    def add_slide(self, layout_index: int = 0):
        """
        Add a new slide with the specified layout.

        Args:
            layout_index: Index of the layout to use (default: 0 - title slide)

        Returns:
            The newly created slide
        """
        layout = self.get_layout_by_index(layout_index)
        return self.prs.slides.add_slide(layout)

    def add_title_slide(self, title: str, subtitle: str = ""):
        """
        Add a title slide.

        Args:
            title: Main title text
            subtitle: Subtitle text

        Returns:
            The newly created slide
        """
        slide = self.add_slide(0)

        if slide.shapes.title:
            slide.shapes.title.text = title

        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        return slide

    def add_text_box(self, slide, text: str, left: float, top: float,
                     width: float, height: float, font_size: int = 18,
                     bold: bool = False, color: Optional[tuple] = None):
        """
        Add a text box to a slide.

        Args:
            slide: The slide to add the text box to
            text: Text content
            left, top: Position in inches
            width, height: Size in inches
            font_size: Font size in points
            bold: Whether text should be bold
            color: RGB color tuple (r, g, b)

        Returns:
            The text box shape
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

        text_frame = textbox.text_frame
        text_frame.text = text

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if color:
                    run.font.color.rgb = RGBColor(*color)

        return textbox

    def add_bullet_points(self, slide, points: List[str], left: float = 1,
                         top: float = 2, width: float = 8, height: float = 4,
                         font_size: int = 18):
        """
        Add bullet points to a slide.

        Args:
            slide: The slide to add bullets to
            points: List of bullet point texts
            left, top: Position in inches
            width, height: Size in inches
            font_size: Font size in points

        Returns:
            The text box shape with bullets
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(points):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            p.text = point
            p.level = 0
            p.font.size = Pt(font_size)

        return textbox

    def format_shape(self, shape, fill_color: Optional[tuple] = None,
                    line_color: Optional[tuple] = None, line_width: float = 1):
        """
        Format a shape with colors and line properties.

        Args:
            shape: The shape to format
            fill_color: RGB tuple for fill color
            line_color: RGB tuple for line color
            line_width: Line width in points
        """
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)

        if line_color:
            shape.line.color.rgb = RGBColor(*line_color)
            shape.line.width = Pt(line_width)

    def save(self, output_path: Path):
        """
        Save the presentation to a file.

        Args:
            output_path: Path where the presentation should be saved
        """
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))

    def get_slide_count(self) -> int:
        """Get the total number of slides in the presentation."""
        return len(self.prs.slides)

    def delete_slide(self, index: int):
        """Delete a slide by index."""
        if 0 <= index < len(self.prs.slides):
            rId = self.prs.slides._sldIdLst[index].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[index]

    def list_layouts(self) -> List[str]:
        """List all available slide layouts."""
        return [layout.name for layout in self.prs.slide_layouts]

    # ========== SPEAKER NOTES ==========

    def add_speaker_notes(self, slide, notes_text: str):
        """
        Add speaker notes to a slide.

        Args:
            slide: The slide to add notes to
            notes_text: The speaker notes text
        """
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    def get_speaker_notes(self, slide) -> str:
        """
        Get speaker notes from a slide.

        Args:
            slide: The slide to get notes from

        Returns:
            The speaker notes text
        """
        notes_slide = slide.notes_slide
        return notes_slide.notes_text_frame.text

    # ========== HYPERLINKS ==========

    def add_hyperlink_to_shape(self, shape, url: str):
        """
        Add a hyperlink to a shape.

        Args:
            shape: The shape to add hyperlink to
            url: The URL to link to
        """
        if hasattr(shape, 'click_action'):
            shape.click_action.hyperlink.address = url

    def add_hyperlink_to_text(self, run, url: str):
        """
        Add a hyperlink to a text run.

        Args:
            run: The text run to add hyperlink to
            url: The URL to link to
        """
        # Create hyperlink
        rPr = run._r.get_or_add_rPr()
        hlinkClick = rPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick')

        if hlinkClick is None:
            hlinkClick = OxmlElement('a:hlinkClick')
            hlinkClick.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
            rPr.append(hlinkClick)

        # Add relationship
        slide_part = run.part
        rId = slide_part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
        hlinkClick.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', rId)

    def add_internal_hyperlink(self, shape, target_slide_index: int):
        """
        Add an internal hyperlink to jump to another slide.

        Args:
            shape: The shape to add hyperlink to
            target_slide_index: Index of the slide to jump to (0-based)
        """
        if hasattr(shape, 'click_action'):
            target_slide = self.prs.slides[target_slide_index]
            shape.click_action.target_slide = target_slide

    # ========== RICH TEXT FORMATTING ==========

    def add_formatted_text_box(self, slide, text_runs: List[Dict[str, Any]],
                               left: float, top: float, width: float, height: float):
        """
        Add a text box with multiple formatted text runs.

        Args:
            slide: The slide to add the text box to
            text_runs: List of dicts with 'text' and optional formatting:
                      {'text': str, 'bold': bool, 'italic': bool, 'underline': bool,
                       'font_size': int, 'color': tuple, 'font_name': str}
            left, top: Position in inches
            width, height: Size in inches

        Returns:
            The text box shape
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

        text_frame = textbox.text_frame
        text_frame.clear()  # Clear default paragraph

        p = text_frame.paragraphs[0]

        for run_spec in text_runs:
            run = p.add_run()
            run.text = run_spec.get('text', '')

            # Apply formatting
            if 'bold' in run_spec:
                run.font.bold = run_spec['bold']
            if 'italic' in run_spec:
                run.font.italic = run_spec['italic']
            if 'underline' in run_spec:
                run.font.underline = run_spec['underline']
            if 'font_size' in run_spec:
                run.font.size = Pt(run_spec['font_size'])
            if 'color' in run_spec:
                run.font.color.rgb = RGBColor(*run_spec['color'])
            if 'font_name' in run_spec:
                run.font.name = run_spec['font_name']

        return textbox

    # ========== SLIDE MANAGEMENT ==========

    def duplicate_slide(self, slide_index: int):
        """
        Duplicate a slide.

        Args:
            slide_index: Index of the slide to duplicate (0-based)

        Returns:
            The newly created duplicate slide
        """
        source_slide = self.prs.slides[slide_index]

        # Get the layout of the source slide
        slide_layout = source_slide.slide_layout

        # Create new slide with same layout
        new_slide = self.prs.slides.add_slide(slide_layout)

        # Copy all shapes from source to new slide
        for shape in source_slide.shapes:
            el = shape.element
            newel = type(el).new()
            newel[:] = el[:]
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        return new_slide

    def reorder_slides(self, new_order: List[int]):
        """
        Reorder slides according to the provided list of indices.

        Args:
            new_order: List of slide indices in the desired order
        """
        xml_slides = self.prs.slides._sldIdLst
        slides = list(xml_slides)

        for idx, slide_idx in enumerate(new_order):
            xml_slides.insert(idx, slides[slide_idx])

    def hide_slide(self, slide_index: int):
        """
        Hide a slide (it remains in the deck but won't show in presentation).

        Args:
            slide_index: Index of the slide to hide (0-based)
        """
        slide = self.prs.slides[slide_index]
        slide_id = self.prs.slides._sldIdLst[slide_index]

        # Set the show attribute to 0 (hidden)
        slide_id.set('show', '0')

    def show_slide(self, slide_index: int):
        """
        Unhide a slide.

        Args:
            slide_index: Index of the slide to show (0-based)
        """
        slide_id = self.prs.slides._sldIdLst[slide_index]

        # Remove the show attribute or set to 1
        if 'show' in slide_id.attrib:
            del slide_id.attrib['show']

    # ========== FOOTER AND SLIDE NUMBERS ==========

    def set_slide_footer(self, slide, footer_text: str = "",
                        show_slide_number: bool = False, show_date: bool = False):
        """
        Set footer, slide number, and date visibility for a slide.

        Args:
            slide: The slide to modify
            footer_text: Footer text to display
            show_slide_number: Whether to show slide number
            show_date: Whether to show date
        """
        # Access the slide's header/footer properties
        if hasattr(slide, 'has_notes_slide'):
            # Try to access placeholders
            for shape in slide.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == 4:  # Footer placeholder
                        shape.text = footer_text
                    elif phf.type == 12 and show_slide_number:  # Slide number
                        shape.text = str(slide.slide_id)
                    elif phf.type == 5 and show_date:  # Date placeholder
                        # Date is typically auto-updated by PowerPoint
                        pass

    def set_presentation_footer(self, footer_text: str = "",
                               show_slide_number: bool = True,
                               show_date: bool = False):
        """
        Set footer settings for all slides in the presentation.

        Args:
            footer_text: Footer text to display
            show_slide_number: Whether to show slide numbers
            show_date: Whether to show date
        """
        for slide in self.prs.slides:
            self.set_slide_footer(slide, footer_text, show_slide_number, show_date)

    # ========== SLIDE DIMENSIONS ==========

    def set_slide_size(self, size_type: str = "widescreen"):
        """
        Set the slide size/dimensions.

        Args:
            size_type: One of 'widescreen' (16:9), 'standard' (4:3), or 'custom'
        """
        if size_type == "widescreen":
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        elif size_type == "standard":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)

        # Update cached dimensions
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def set_custom_slide_size(self, width_inches: float, height_inches: float):
        """
        Set custom slide dimensions.

        Args:
            width_inches: Width in inches
            height_inches: Height in inches
        """
        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    # ========== IMAGE ENHANCEMENTS ==========

    def add_image_with_alt_text(self, slide, image_path: Path, left: float, top: float,
                                width: Optional[float] = None, height: Optional[float] = None,
                                alt_text: str = ""):
        """
        Add an image with alt text for accessibility.

        Args:
            slide: The slide to add the image to
            image_path: Path to the image file
            left, top: Position in inches
            width, height: Optional size in inches
            alt_text: Alternative text for accessibility

        Returns:
            The image shape
        """
        if width and height:
            pic = slide.shapes.add_picture(
                str(image_path), Inches(left), Inches(top),
                width=Inches(width), height=Inches(height)
            )
        else:
            pic = slide.shapes.add_picture(
                str(image_path), Inches(left), Inches(top)
            )

        # Set alt text
        pic._element.nvPicPr.cNvPr.set('descr', alt_text)

        return pic

    def set_image_transparency(self, image_shape, transparency: float):
        """
        Set image transparency (alpha).

        Args:
            image_shape: The image shape
            transparency: Transparency value from 0.0 (opaque) to 1.0 (fully transparent)
        """
        # This requires accessing the XML directly
        # Convert transparency to percentage (0-100000)
        alpha = int((1 - transparency) * 100000)

        # Access the fill properties
        pic = image_shape._element
        blipFill = pic.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')
        if blipFill is not None:
            blip = blipFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip is not None:
                # Add alpha mod fix
                alphaModFix = OxmlElement('a:alphaModFix')
                alphaModFix.set('amt', str(alpha))
                blip.append(alphaModFix)

    # ========== GROUPING AND LAYERING ==========

    def group_shapes(self, slide, shape_indices: List[int]):
        """
        Group multiple shapes together.

        Note: This is a simplified implementation. Full grouping support
        in python-pptx is limited.

        Args:
            slide: The slide containing the shapes
            shape_indices: List of shape indices to group

        Returns:
            The group shape (if supported)
        """
        # python-pptx has limited grouping support
        # This is a placeholder for future enhancement
        shapes_to_group = [slide.shapes[i] for i in shape_indices]
        # Full implementation would require XML manipulation
        return shapes_to_group

    def send_to_back(self, slide, shape_index: int):
        """
        Send a shape to the back (bottom of z-order).

        Args:
            slide: The slide containing the shape
            shape_index: Index of the shape to move
        """
        shape = slide.shapes[shape_index]
        spTree = slide._element.cSld.spTree

        # Move shape element to beginning of spTree (after nvGrpSpPr and grpSpPr)
        spTree.remove(shape._element)
        spTree.insert(2, shape._element)

    def bring_to_front(self, slide, shape_index: int):
        """
        Bring a shape to the front (top of z-order).

        Args:
            slide: The slide containing the shape
            shape_index: Index of the shape to move
        """
        shape = slide.shapes[shape_index]
        spTree = slide._element.cSld.spTree

        # Move shape element to end of spTree
        spTree.remove(shape._element)
        spTree.append(shape._element)
