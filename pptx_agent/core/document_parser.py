"""
Document Parser - Extract text from various document formats.

Supports:
- .txt, .md - Plain text files
- .docx - Microsoft Word documents
- .pptx - Microsoft PowerPoint presentations
- .xlsx - Microsoft Excel spreadsheets
"""

from pathlib import Path
from typing import Optional, Dict, Any, List
import logging

logger = logging.getLogger(__name__)


class DocumentParser:
    """Parse various document formats and extract text content."""

    @staticmethod
    def parse_file(file_path: Path) -> str:
        """
        Parse a document file and extract its text content.

        Args:
            file_path: Path to the document file

        Returns:
            Extracted text content as a string

        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Get file extension
        ext = file_path.suffix.lower()

        # Route to appropriate parser
        if ext in ['.txt', '.md']:
            return DocumentParser._parse_text_file(file_path)
        elif ext == '.docx':
            return DocumentParser._parse_docx(file_path)
        elif ext == '.pptx':
            return DocumentParser._parse_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return DocumentParser._parse_excel(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    @staticmethod
    def _parse_text_file(file_path: Path) -> str:
        """Parse plain text files (.txt, .md)."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()

    @staticmethod
    def _parse_docx(file_path: Path) -> str:
        """
        Parse Microsoft Word documents (.docx).

        Extracts:
        - Paragraph text
        - Table content
        - Headers and footers
        """
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return f"[Error: python-docx not installed. Cannot parse {file_path.name}]"

        try:
            doc = Document(str(file_path))
            text_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table in doc.tables:
                text_parts.append("\n[TABLE]")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
                text_parts.append("[/TABLE]\n")

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error parsing DOCX file {file_path}: {e}")
            return f"[Error parsing {file_path.name}: {str(e)}]"

    @staticmethod
    def _parse_pptx(file_path: Path) -> str:
        """
        Parse Microsoft PowerPoint presentations (.pptx).

        Extracts:
        - Slide titles
        - Text from all shapes
        - Table content
        - Speaker notes
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return f"[Error: python-pptx not installed. Cannot parse {file_path.name}]"

        try:
            prs = Presentation(str(file_path))
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n[SLIDE {i}]")

                # Extract title
                if slide.shapes.title:
                    text_parts.append(f"Title: {slide.shapes.title.text}")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape != slide.shapes.title:  # Skip title (already added)
                            text_parts.append(shape.text)

                    # Extract table content
                    if shape.shape_type == 19:  # Table
                        try:
                            table = shape.table
                            text_parts.append("\n[TABLE]")
                            for row in table.rows:
                                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                                if row_text.strip():
                                    text_parts.append(row_text)
                            text_parts.append("[/TABLE]\n")
                        except Exception:
                            pass

                # Extract speaker notes
                try:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        text_parts.append(f"Speaker Notes: {notes_text}")
                except Exception:
                    pass

                text_parts.append("[/SLIDE]\n")

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {e}")
            return f"[Error parsing {file_path.name}: {str(e)}]"

    @staticmethod
    def _parse_excel(file_path: Path) -> str:
        """
        Parse Microsoft Excel spreadsheets (.xlsx, .xls).

        Extracts:
        - Sheet names
        - Cell values from all sheets
        - Formatted as structured text
        """
        try:
            import pandas as pd
        except ImportError:
            logger.error("pandas not installed. Install with: pip install pandas openpyxl")
            return f"[Error: pandas not installed. Cannot parse {file_path.name}]"

        try:
            # Read all sheets
            excel_file = pd.ExcelFile(str(file_path))
            text_parts = []

            for sheet_name in excel_file.sheet_names:
                text_parts.append(f"\n[SHEET: {sheet_name}]")

                # Read sheet
                df = pd.read_excel(excel_file, sheet_name=sheet_name)

                # Convert to string representation
                if not df.empty:
                    # Get column headers
                    headers = " | ".join(str(col) for col in df.columns)
                    text_parts.append(headers)
                    text_parts.append("-" * len(headers))

                    # Get rows (limit to first 100 rows to avoid huge text)
                    for idx, row in df.head(100).iterrows():
                        row_text = " | ".join(str(val) for val in row.values)
                        text_parts.append(row_text)

                    if len(df) > 100:
                        text_parts.append(f"... ({len(df) - 100} more rows)")

                text_parts.append("[/SHEET]\n")

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error parsing Excel file {file_path}: {e}")
            return f"[Error parsing {file_path.name}: {str(e)}]"

    @staticmethod
    def parse_multiple_files(file_paths: List[Path]) -> str:
        """
        Parse multiple document files and combine their content.

        Args:
            file_paths: List of paths to document files

        Returns:
            Combined text content from all files
        """
        combined_text = []

        for file_path in file_paths:
            combined_text.append(f"\n{'='*60}")
            combined_text.append(f"SOURCE: {file_path.name}")
            combined_text.append('='*60 + "\n")

            try:
                content = DocumentParser.parse_file(file_path)
                combined_text.append(content)
            except Exception as e:
                logger.error(f"Error parsing {file_path}: {e}")
                combined_text.append(f"[Error: Could not parse {file_path.name}]")

        return "\n".join(combined_text)

    @staticmethod
    def get_document_summary(file_path: Path) -> Dict[str, Any]:
        """
        Get a summary of document properties.

        Args:
            file_path: Path to the document file

        Returns:
            Dictionary with document metadata
        """
        summary = {
            "filename": file_path.name,
            "extension": file_path.suffix.lower(),
            "size_bytes": file_path.stat().st_size,
            "supported": file_path.suffix.lower() in ['.txt', '.md', '.docx', '.pptx', '.xlsx', '.xls']
        }

        # Add format-specific info
        if summary["extension"] == ".docx":
            try:
                from docx import Document
                doc = Document(str(file_path))
                summary["paragraphs"] = len(doc.paragraphs)
                summary["tables"] = len(doc.tables)
            except Exception:
                pass

        elif summary["extension"] == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                summary["slides"] = len(prs.slides)
            except Exception:
                pass

        elif summary["extension"] in [".xlsx", ".xls"]:
            try:
                import pandas as pd
                excel_file = pd.ExcelFile(str(file_path))
                summary["sheets"] = len(excel_file.sheet_names)
                summary["sheet_names"] = excel_file.sheet_names
            except Exception:
                pass

        return summary
