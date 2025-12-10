"""
SmartArt Builder - Creates SmartArt-like diagrams and process flows.

Note: python-pptx doesn't have native SmartArt support, so we create
SmartArt-like visuals using shapes and text boxes.
"""

from typing import List, Dict, Any, Optional, Tuple
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class SmartArtBuilder:
    """Builds SmartArt-like diagrams using shapes and text."""

    @staticmethod
    def add_process_flow(handler, title: str, steps: List[str],
                        layout_index: int = 1,
                        colors: Optional[List[Tuple[int, int, int]]] = None):
        """
        Create a horizontal process flow diagram.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            steps: List of process steps
            layout_index: Layout to use
            colors: Optional list of RGB color tuples

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Default colors (blue gradient)
        if not colors:
            colors = [(68, 114, 196), (112, 173, 71), (237, 125, 49)]

        num_steps = len(steps)
        box_width = 7.5 / num_steps
        box_height = 1.5
        start_left = 1.0
        top = 2.5

        for i, step in enumerate(steps):
            left = start_left + (i * (box_width + 0.2))
            color = colors[i % len(colors)]

            # Add rounded rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top),
                Inches(box_width), Inches(box_height)
            )

            # Style the shape
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*color)
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)

            # Add text
            text_frame = shape.text_frame
            text_frame.text = step
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_SHAPE.MIDDLE

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True

            # Add arrow between steps
            if i < num_steps - 1:
                arrow_left = left + box_width + 0.05
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(arrow_left), Inches(top + 0.5),
                    Inches(0.3), Inches(0.5)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
                arrow.line.fill.background()

        return slide

    @staticmethod
    def add_cycle_diagram(handler, title: str, items: List[str],
                         layout_index: int = 1):
        """
        Create a circular cycle diagram.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            items: List of items in the cycle
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_items = len(items)
        center_x, center_y = 5.0, 4.0
        radius = 2.0
        box_size = 1.2

        import math

        for i, item in enumerate(items):
            # Calculate position in circle
            angle = (2 * math.pi / num_items) * i - (math.pi / 2)
            x = center_x + radius * math.cos(angle) - box_size / 2
            y = center_y + radius * math.sin(angle) - box_size / 2

            # Add oval shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(box_size), Inches(box_size)
            )

            # Style
            color_idx = i % 3
            colors = [(68, 114, 196), (112, 173, 71), (237, 125, 49)]
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*colors[color_idx])
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(3)

            # Add text
            text_frame = shape.text_frame
            text_frame.text = item
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_SHAPE.MIDDLE

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True

        return slide

    @staticmethod
    def add_hierarchy_diagram(handler, title: str,
                             root: str, children: List[str],
                             layout_index: int = 1):
        """
        Create a simple hierarchy diagram.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            root: Root/top level item
            children: List of child items
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Root box
        root_width, root_height = 3.0, 1.0
        root_left = 3.5
        root_top = 2.0

        root_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(root_left), Inches(root_top),
            Inches(root_width), Inches(root_height)
        )

        root_shape.fill.solid()
        root_shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
        root_shape.line.color.rgb = RGBColor(255, 255, 255)

        text_frame = root_shape.text_frame
        text_frame.text = root
        text_frame.vertical_anchor = MSO_SHAPE.MIDDLE
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True

        # Child boxes
        num_children = len(children)
        child_width = 7.0 / num_children if num_children > 0 else 2.0
        child_height = 0.8
        child_top = 4.0
        start_left = 1.5

        for i, child in enumerate(children):
            left = start_left + (i * child_width)

            child_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(child_top),
                Inches(child_width - 0.2), Inches(child_height)
            )

            child_shape.fill.solid()
            child_shape.fill.fore_color.rgb = RGBColor(112, 173, 71)
            child_shape.line.color.rgb = RGBColor(255, 255, 255)

            text_frame = child_shape.text_frame
            text_frame.text = child
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_SHAPE.MIDDLE

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

            # Add connecting line
            line = slide.shapes.add_connector(
                1,  # Straight connector
                Inches(root_left + root_width / 2), Inches(root_top + root_height),
                Inches(left + (child_width - 0.2) / 2), Inches(child_top)
            )
            line.line.color.rgb = RGBColor(100, 100, 100)

        return slide

    @staticmethod
    def add_comparison_diagram(handler, title: str,
                              left_items: List[str], right_items: List[str],
                              left_label: str = "Option A",
                              right_label: str = "Option B",
                              layout_index: int = 1):
        """
        Create a comparison diagram (two columns).

        Args:
            handler: PPTXHandler instance
            title: Slide title
            left_items: Items for left column
            right_items: Items for right column
            left_label: Label for left column
            right_label: Label for right column
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Left column header
        left_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(2.0),
            Inches(4.0), Inches(0.6)
        )
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = RGBColor(68, 114, 196)
        left_header.text_frame.text = left_label
        left_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        left_header.text_frame.paragraphs[0].font.bold = True
        left_header.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Right column header
        right_header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(2.0),
            Inches(4.0), Inches(0.6)
        )
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = RGBColor(237, 125, 49)
        right_header.text_frame.text = right_label
        right_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        right_header.text_frame.paragraphs[0].font.bold = True
        right_header.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Left items
        current_top = 2.8
        for item in left_items:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(current_top),
                Inches(3.6), Inches(0.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(217, 225, 242)
            shape.text_frame.text = item
            shape.text_frame.paragraphs[0].font.size = Pt(12)
            current_top += 0.6

        # Right items
        current_top = 2.8
        for item in right_items:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.7), Inches(current_top),
                Inches(3.6), Inches(0.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(252, 228, 214)
            shape.text_frame.text = item
            shape.text_frame.paragraphs[0].font.size = Pt(12)
            current_top += 0.6

        return slide

    @staticmethod
    def add_venn_diagram(handler, title: str,
                        left_label: str, right_label: str,
                        left_items: List[str], right_items: List[str],
                        overlap_items: List[str],
                        layout_index: int = 1):
        """
        Create a Venn diagram using circles.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            left_label: Label for left circle
            right_label: Label for right circle
            left_items: Items unique to left
            right_items: Items unique to right
            overlap_items: Items in both
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Left circle
        left_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.5), Inches(2.0),
            Inches(3.5), Inches(3.5)
        )
        left_circle.fill.solid()
        left_circle.fill.fore_color.rgb = RGBColor(68, 114, 196)
        left_circle.fill.transparency = 0.5

        # Right circle
        right_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.0), Inches(2.0),
            Inches(3.5), Inches(3.5)
        )
        right_circle.fill.solid()
        right_circle.fill.fore_color.rgb = RGBColor(237, 125, 49)
        right_circle.fill.transparency = 0.5

        # Labels
        handler.add_text_box(slide, left_label, 2.0, 1.5, 2.0, 0.4, font_size=14, bold=True)
        handler.add_text_box(slide, right_label, 6.5, 1.5, 2.0, 0.4, font_size=14, bold=True)

        # Add text for unique items (simplified - actual positioning would be more complex)
        left_text = "\n".join(left_items)
        handler.add_text_box(slide, left_text, 1.7, 3.0, 1.5, 2.0, font_size=10)

        right_text = "\n".join(right_items)
        handler.add_text_box(slide, right_text, 6.5, 3.0, 1.5, 2.0, font_size=10)

        overlap_text = "\n".join(overlap_items)
        handler.add_text_box(slide, overlap_text, 4.0, 3.5, 2.0, 1.5, font_size=10)

        return slide

    @staticmethod
    def add_timeline(handler, title: str, events: List[Dict[str, str]],
                    layout_index: int = 1):
        """
        Create a horizontal timeline.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            events: List of dicts with 'date' and 'event' keys
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_events = len(events)
        timeline_top = 3.5
        timeline_left = 1.0
        timeline_width = 8.0

        # Draw main timeline line
        line = slide.shapes.add_connector(
            1,  # Straight
            Inches(timeline_left), Inches(timeline_top),
            Inches(timeline_left + timeline_width), Inches(timeline_top)
        )
        line.line.color.rgb = RGBColor(100, 100, 100)
        line.line.width = Pt(3)

        # Add events
        spacing = timeline_width / (num_events - 1) if num_events > 1 else 0

        for i, event in enumerate(events):
            x = timeline_left + (i * spacing)

            # Marker circle
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.15), Inches(timeline_top - 0.15),
                Inches(0.3), Inches(0.3)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(68, 114, 196)

            # Date (above timeline)
            handler.add_text_box(
                slide, event.get('date', ''),
                x - 0.5, timeline_top - 0.8,
                1.0, 0.4,
                font_size=11, bold=True
            )

            # Event (below timeline)
            handler.add_text_box(
                slide, event.get('event', ''),
                x - 0.5, timeline_top + 0.3,
                1.0, 0.8,
                font_size=10
            )

        return slide
