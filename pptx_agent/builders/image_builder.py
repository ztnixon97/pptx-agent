"""
Image Slide Builder - Creates slides with images.
"""

from typing import Optional
from pathlib import Path
from PIL import Image
from pptx.util import Inches


class ImageSlideBuilder:
    """Builds slides with images."""

    @staticmethod
    def add_image_slide(handler, title: str, image_path: Path,
                       layout_index: int = 1,
                       left: float = 1.5, top: float = 2,
                       width: Optional[float] = None,
                       height: Optional[float] = None,
                       caption: str = ""):
        """
        Add a slide with an image.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            image_path: Path to the image file
            layout_index: Layout to use
            left, top: Position in inches
            width: Width in inches (auto if None)
            height: Height in inches (auto if None)
            caption: Optional image caption

        Returns:
            The created slide
        """
        if not image_path.exists():
            raise FileNotFoundError(f"Image not found: {image_path}")

        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Add image
        kwargs = {'left': Inches(left), 'top': Inches(top)}
        if width:
            kwargs['width'] = Inches(width)
        if height:
            kwargs['height'] = Inches(height)

        pic = slide.shapes.add_picture(str(image_path), **kwargs)

        # Add caption if provided
        if caption:
            caption_top = top + (pic.height / 914400)  # Convert EMUs to inches
            handler.add_text_box(
                slide, caption,
                left=left, top=caption_top + 0.2,
                width=pic.width / 914400, height=0.4,
                font_size=12
            )

        return slide

    @staticmethod
    def add_full_image_slide(handler, image_path: Path,
                            layout_index: int = 6):
        """
        Add a slide with a full-screen image (no title).

        Args:
            handler: PPTXHandler instance
            image_path: Path to the image file
            layout_index: Layout to use (default: blank layout)

        Returns:
            The created slide
        """
        if not image_path.exists():
            raise FileNotFoundError(f"Image not found: {image_path}")

        slide = handler.add_slide(layout_index)

        # Calculate dimensions to fill slide
        img = Image.open(image_path)
        img_width, img_height = img.size
        img_ratio = img_width / img_height

        slide_width = handler.slide_width / 914400  # Convert to inches
        slide_height = handler.slide_height / 914400

        slide_ratio = slide_width / slide_height

        if img_ratio > slide_ratio:
            # Image is wider, fit to width
            width = slide_width
            height = slide_width / img_ratio
            left = 0
            top = (slide_height - height) / 2
        else:
            # Image is taller, fit to height
            height = slide_height
            width = slide_height * img_ratio
            top = 0
            left = (slide_width - width) / 2

        slide.shapes.add_picture(
            str(image_path),
            Inches(left), Inches(top),
            width=Inches(width), height=Inches(height)
        )

        return slide

    @staticmethod
    def add_multiple_images_slide(handler, title: str,
                                 image_paths: list[Path],
                                 layout_index: int = 1):
        """
        Add a slide with multiple images arranged in a grid.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            image_paths: List of paths to image files
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_images = len(image_paths)
        if num_images == 0:
            return slide

        # Determine grid layout
        if num_images == 1:
            cols, rows = 1, 1
        elif num_images == 2:
            cols, rows = 2, 1
        elif num_images <= 4:
            cols, rows = 2, 2
        else:
            cols, rows = 3, 2

        # Calculate dimensions
        start_left = 1.0
        start_top = 2.0
        available_width = 8.0
        available_height = 4.5

        img_width = available_width / cols - 0.2
        img_height = available_height / rows - 0.2

        # Add images
        for idx, img_path in enumerate(image_paths[:6]):  # Max 6 images
            if not img_path.exists():
                continue

            row = idx // cols
            col = idx % cols

            left = start_left + col * (img_width + 0.2)
            top = start_top + row * (img_height + 0.2)

            try:
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(left), Inches(top),
                    width=Inches(img_width)
                )
            except Exception:
                continue

        return slide

    @staticmethod
    def add_image_with_text_slide(handler, title: str, image_path: Path,
                                  text: str, image_on_left: bool = True,
                                  layout_index: int = 1):
        """
        Add a slide with an image and text side by side.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            image_path: Path to the image file
            text: Text content
            image_on_left: If True, image on left, text on right; otherwise reversed
            layout_index: Layout to use

        Returns:
            The created slide
        """
        if not image_path.exists():
            raise FileNotFoundError(f"Image not found: {image_path}")

        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        if image_on_left:
            img_left, img_top = 0.5, 2.0
            text_left, text_top = 5.0, 2.0
        else:
            img_left, img_top = 5.0, 2.0
            text_left, text_top = 0.5, 2.0

        # Add image
        slide.shapes.add_picture(
            str(image_path),
            Inches(img_left), Inches(img_top),
            width=Inches(4.0)
        )

        # Add text
        handler.add_text_box(
            slide, text,
            left=text_left, top=text_top,
            width=4.5, height=4.5,
            font_size=14
        )

        return slide

    @staticmethod
    def validate_image(image_path: Path) -> bool:
        """
        Validate that an image file is readable.

        Args:
            image_path: Path to the image file

        Returns:
            True if valid, False otherwise
        """
        try:
            if not image_path.exists():
                return False
            img = Image.open(image_path)
            img.verify()
            return True
        except Exception:
            return False
