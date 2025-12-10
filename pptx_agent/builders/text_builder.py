"""
Text Slide Builder - Creates text-based slides.
"""

from typing import List, Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class TextSlideBuilder:
    """Builds text-based slides with various formatting options."""

    @staticmethod
    def add_content_slide(handler, title: str, content: str,
                         layout_index: int = 1):
        """
        Add a basic content slide with title and text.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            content: Main content text
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Try to use content placeholder
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = content
        else:
            # Add text box if no placeholder
            handler.add_text_box(slide, content, 1, 2, 8, 4)

        return slide

    @staticmethod
    def add_bullet_slide(handler, title: str, bullet_points: List[str],
                        layout_index: int = 1, font_size: int = 18):
        """
        Add a slide with bullet points.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            bullet_points: List of bullet point texts
            layout_index: Layout to use
            font_size: Font size for bullets

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Add bullet points
        handler.add_bullet_points(
            slide, bullet_points,
            left=1, top=2, width=8, height=4.5,
            font_size=font_size
        )

        return slide

    @staticmethod
    def add_two_column_slide(handler, title: str,
                           left_content: str, right_content: str,
                           layout_index: int = 1):
        """
        Add a slide with two columns of text.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            left_content: Content for left column
            right_content: Content for right column
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Left column
        handler.add_text_box(
            slide, left_content,
            left=0.5, top=2, width=4.25, height=4.5,
            font_size=16
        )

        # Right column
        handler.add_text_box(
            slide, right_content,
            left=5.25, top=2, width=4.25, height=4.5,
            font_size=16
        )

        return slide

    @staticmethod
    def add_section_slide(handler, section_title: str,
                         subtitle: str = "", layout_index: int = 2):
        """
        Add a section divider slide.

        Args:
            handler: PPTXHandler instance
            section_title: Section title
            subtitle: Optional subtitle
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = section_title

        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        return slide

    @staticmethod
    def add_quote_slide(handler, quote: str, author: str = "",
                       layout_index: int = 6):
        """
        Add a slide with a quote.

        Args:
            handler: PPTXHandler instance
            quote: The quote text
            author: Author of the quote
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        # Add quote
        quote_box = handler.add_text_box(
            slide, f'"{quote}"',
            left=1.5, top=2, width=7, height=3,
            font_size=24
        )

        # Center align and italicize
        for paragraph in quote_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.italic = True

        # Add author if provided
        if author:
            author_box = handler.add_text_box(
                slide, f"— {author}",
                left=1.5, top=5.2, width=7, height=0.5,
                font_size=18
            )
            for paragraph in author_box.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.RIGHT

        return slide

    @staticmethod
    def add_blank_slide(handler, layout_index: int = 6):
        """
        Add a blank slide for custom content.

        Args:
            handler: PPTXHandler instance
            layout_index: Layout to use

        Returns:
            The created slide
        """
        return handler.add_slide(layout_index)

    @staticmethod
    def add_formatted_text_slide(handler, title: str, paragraphs: List[Dict[str, Any]],
                                layout_index: int = 1):
        """
        Add a slide with formatted text paragraphs.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            paragraphs: List of paragraph dictionaries with 'text', 'font_size', 'bold', 'color'
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        current_top = 2.0
        for para in paragraphs:
            text = para.get('text', '')
            font_size = para.get('font_size', 16)
            bold = para.get('bold', False)
            color = para.get('color', None)

            # Estimate height based on text length
            height = max(0.5, len(text) / 100)

            handler.add_text_box(
                slide, text,
                left=1, top=current_top, width=8, height=height,
                font_size=font_size, bold=bold, color=color
            )

            current_top += height + 0.2

        return slide
