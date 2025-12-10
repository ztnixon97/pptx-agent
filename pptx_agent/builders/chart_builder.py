"""
Chart Slide Builder - Creates slides with charts and graphs.
"""

import io
from typing import List, Optional, Dict, Any
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend

from pptx.util import Inches, Pt


class ChartSlideBuilder:
    """Builds slides with charts and graphs using matplotlib."""

    @staticmethod
    def _create_bar_chart(categories: List[str], series: List[Dict[str, Any]],
                         title: str = "", xlabel: str = "", ylabel: str = ""):
        """Create a bar chart."""
        fig, ax = plt.subplots(figsize=(10, 6))

        x = range(len(categories))
        width = 0.8 / len(series)

        for i, s in enumerate(series):
            offset = width * i - (width * len(series) / 2 - width / 2)
            ax.bar([pos + offset for pos in x], s['values'],
                  width, label=s['name'])

        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel)
        ax.set_title(title)
        ax.set_xticks(x)
        ax.set_xticklabels(categories)
        ax.legend()
        ax.grid(axis='y', alpha=0.3)

        return fig

    @staticmethod
    def _create_line_chart(categories: List[str], series: List[Dict[str, Any]],
                          title: str = "", xlabel: str = "", ylabel: str = ""):
        """Create a line chart."""
        fig, ax = plt.subplots(figsize=(10, 6))

        for s in series:
            ax.plot(categories, s['values'], marker='o', label=s['name'])

        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel)
        ax.set_title(title)
        ax.legend()
        ax.grid(True, alpha=0.3)

        return fig

    @staticmethod
    def _create_pie_chart(categories: List[str], values: List[float],
                         title: str = ""):
        """Create a pie chart."""
        fig, ax = plt.subplots(figsize=(10, 6))

        ax.pie(values, labels=categories, autopct='%1.1f%%', startangle=90)
        ax.set_title(title)
        ax.axis('equal')

        return fig

    @staticmethod
    def _fig_to_image(fig) -> io.BytesIO:
        """Convert matplotlib figure to image bytes."""
        img_buffer = io.BytesIO()
        fig.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close(fig)
        return img_buffer

    @staticmethod
    def add_bar_chart_slide(handler, title: str, categories: List[str],
                           series: List[Dict[str, Any]],
                           layout_index: int = 1,
                           chart_title: str = "",
                           xlabel: str = "", ylabel: str = ""):
        """
        Add a slide with a bar chart.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            categories: List of category labels
            series: List of series dictionaries with 'name' and 'values'
            layout_index: Layout to use
            chart_title: Title for the chart
            xlabel: X-axis label
            ylabel: Y-axis label

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Create chart
        fig = ChartSlideBuilder._create_bar_chart(
            categories, series, chart_title, xlabel, ylabel
        )
        img_buffer = ChartSlideBuilder._fig_to_image(fig)

        # Add image to slide
        slide.shapes.add_picture(
            img_buffer,
            Inches(1.5), Inches(2),
            width=Inches(7)
        )

        return slide

    @staticmethod
    def add_line_chart_slide(handler, title: str, categories: List[str],
                            series: List[Dict[str, Any]],
                            layout_index: int = 1,
                            chart_title: str = "",
                            xlabel: str = "", ylabel: str = ""):
        """
        Add a slide with a line chart.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            categories: List of category labels
            series: List of series dictionaries with 'name' and 'values'
            layout_index: Layout to use
            chart_title: Title for the chart
            xlabel: X-axis label
            ylabel: Y-axis label

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Create chart
        fig = ChartSlideBuilder._create_line_chart(
            categories, series, chart_title, xlabel, ylabel
        )
        img_buffer = ChartSlideBuilder._fig_to_image(fig)

        # Add image to slide
        slide.shapes.add_picture(
            img_buffer,
            Inches(1.5), Inches(2),
            width=Inches(7)
        )

        return slide

    @staticmethod
    def add_pie_chart_slide(handler, title: str, categories: List[str],
                           values: List[float],
                           layout_index: int = 1,
                           chart_title: str = ""):
        """
        Add a slide with a pie chart.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            categories: List of category labels
            values: List of values for each category
            layout_index: Layout to use
            chart_title: Title for the chart

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Create chart
        fig = ChartSlideBuilder._create_pie_chart(
            categories, values, chart_title
        )
        img_buffer = ChartSlideBuilder._fig_to_image(fig)

        # Add image to slide
        slide.shapes.add_picture(
            img_buffer,
            Inches(2.5), Inches(2),
            width=Inches(5)
        )

        return slide

    @staticmethod
    def add_custom_chart_slide(handler, title: str, fig,
                              layout_index: int = 1,
                              left: float = 1.5, top: float = 2,
                              width: float = 7):
        """
        Add a slide with a custom matplotlib figure.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            fig: Matplotlib figure object
            layout_index: Layout to use
            left, top: Position in inches
            width: Width in inches

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Convert figure to image
        img_buffer = ChartSlideBuilder._fig_to_image(fig)

        # Add image to slide
        slide.shapes.add_picture(
            img_buffer,
            Inches(left), Inches(top),
            width=Inches(width)
        )

        return slide

    @staticmethod
    def create_chart_from_data(chart_type: str, categories: List[str],
                              series: List[Dict[str, Any]],
                              title: str = "", xlabel: str = "",
                              ylabel: str = ""):
        """
        Create a chart figure from data.

        Args:
            chart_type: Type of chart ('bar', 'line', 'pie')
            categories: List of category labels
            series: List of series dictionaries
            title: Chart title
            xlabel: X-axis label
            ylabel: Y-axis label

        Returns:
            Matplotlib figure object
        """
        if chart_type.lower() == 'bar':
            return ChartSlideBuilder._create_bar_chart(
                categories, series, title, xlabel, ylabel
            )
        elif chart_type.lower() == 'line':
            return ChartSlideBuilder._create_line_chart(
                categories, series, title, xlabel, ylabel
            )
        elif chart_type.lower() == 'pie' and series:
            return ChartSlideBuilder._create_pie_chart(
                categories, series[0]['values'], title
            )
        else:
            return ChartSlideBuilder._create_bar_chart(
                categories, series, title, xlabel, ylabel
            )
