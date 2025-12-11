"""
Table Slide Builder - Creates slides with tables.
"""

from typing import List, Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class TableSlideBuilder:
    """Builds slides with tables."""

    @staticmethod
    def add_table_slide(handler, title: str, headers: List[str],
                       rows: List[List[str]], layout_index: int = 1,
                       left: float = 1.0, top: float = 2.0,
                       width: float = 8.0, height: float = 4.0):
        """
        Add a slide with a table.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            headers: List of column headers
            rows: List of rows (each row is a list of cell values)
            layout_index: Layout to use
            left, top: Position in inches
            width, height: Size in inches

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Calculate table dimensions
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row

        # Add table
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        table = table_shape.table

        # Set column headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)  # Blue header
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Fill in data rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_value in enumerate(row_data):
                if col_idx < num_cols:  # Ensure we don't exceed column count
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)
                    cell.text_frame.paragraphs[0].font.size = Pt(11)

                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(217, 225, 242)

        return slide

    @staticmethod
    def add_comparison_table_slide(handler, title: str,
                                   comparison_data: Dict[str, List[str]],
                                   layout_index: int = 1):
        """
        Add a slide with a comparison table.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            comparison_data: Dictionary with column names as keys and data lists as values
            layout_index: Layout to use

        Returns:
            The created slide
        """
        headers = list(comparison_data.keys())
        num_rows = max(len(values) for values in comparison_data.values())

        rows = []
        for i in range(num_rows):
            row = []
            for header in headers:
                values = comparison_data[header]
                row.append(values[i] if i < len(values) else "")
            rows.append(row)

        return TableSlideBuilder.add_table_slide(
            handler, title, headers, rows, layout_index
        )

    @staticmethod
    def add_styled_table_slide(handler, title: str, headers: List[str],
                              rows: List[List[str]],
                              header_color: tuple = (68, 114, 196),
                              alt_row_color: tuple = (217, 225, 242),
                              layout_index: int = 1):
        """
        Add a slide with a styled table with custom colors.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            headers: List of column headers
            rows: List of rows
            header_color: RGB tuple for header background
            alt_row_color: RGB tuple for alternate rows
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_cols = len(headers)
        num_rows = len(rows) + 1

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(1.0), Inches(2.0),
            Inches(8.0), Inches(4.0)
        )

        table = table_shape.table

        # Set headers with custom color
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*header_color)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Fill data with styling
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_value in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)
                    cell.text_frame.paragraphs[0].font.size = Pt(11)

                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*alt_row_color)

        return slide

    @staticmethod
    def add_summary_table_slide(handler, title: str,
                               labels: List[str], values: List[str],
                               layout_index: int = 1):
        """
        Add a slide with a summary table (2 columns: label and value).

        Args:
            handler: PPTXHandler instance
            title: Slide title
            labels: List of labels
            values: List of corresponding values
            layout_index: Layout to use

        Returns:
            The created slide
        """
        rows = [[label, value] for label, value in zip(labels, values)]

        return TableSlideBuilder.add_table_slide(
            handler, title,
            headers=["Item", "Value"],
            rows=rows,
            layout_index=layout_index,
            left=2.0, width=6.0
        )

    @staticmethod
    def add_advanced_table_slide(handler, title: str, headers: List[str],
                                rows: List[List[str]],
                                merge_cells: Optional[List[tuple]] = None,
                                cell_styles: Optional[Dict[tuple, Dict[str, Any]]] = None,
                                layout_index: int = 1):
        """
        Add a slide with an advanced table supporting cell merging and individual cell styling.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            headers: List of column headers
            rows: List of rows
            merge_cells: List of tuples (row1, col1, row2, col2) to merge cells
            cell_styles: Dict mapping (row, col) to style dict:
                        {'fill': (r,g,b), 'bold': bool, 'italic': bool,
                         'font_size': int, 'font_color': (r,g,b), 'align': str}
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row

        # Add table
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(1.0), Inches(2.0),
            Inches(8.0), Inches(4.0)
        )

        table = table_shape.table

        # Set column headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Fill in data rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_value in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)
                    cell.text_frame.paragraphs[0].font.size = Pt(11)

        # Apply cell merging
        if merge_cells:
            for row1, col1, row2, col2 in merge_cells:
                try:
                    table.cell(row1, col1).merge(table.cell(row2, col2))
                except Exception:
                    pass  # Skip if merge fails

        # Apply individual cell styles
        if cell_styles:
            for (row, col), style in cell_styles.items():
                try:
                    cell = table.cell(row, col)

                    # Apply fill color
                    if 'fill' in style:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*style['fill'])

                    # Apply text formatting
                    para = cell.text_frame.paragraphs[0]

                    if 'bold' in style:
                        para.font.bold = style['bold']
                    if 'italic' in style:
                        para.font.italic = style['italic']
                    if 'font_size' in style:
                        para.font.size = Pt(style['font_size'])
                    if 'font_color' in style:
                        para.font.color.rgb = RGBColor(*style['font_color'])
                    if 'align' in style:
                        align_map = {
                            'left': PP_ALIGN.LEFT,
                            'center': PP_ALIGN.CENTER,
                            'right': PP_ALIGN.RIGHT
                        }
                        para.alignment = align_map.get(style['align'], PP_ALIGN.LEFT)

                except Exception:
                    pass  # Skip if styling fails

        return slide
