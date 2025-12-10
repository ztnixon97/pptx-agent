"""
Shapes Builder - Creates custom shapes, callouts, and annotations.
"""

from typing import List, Optional, Tuple
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class ShapesBuilder:
    """Builds slides with custom shapes and annotations."""

    # Common shape types for easy reference
    SHAPES = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'oval': MSO_SHAPE.OVAL,
        'circle': MSO_SHAPE.OVAL,  # Same as oval
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'diamond': MSO_SHAPE.DIAMOND,
        'pentagon': MSO_SHAPE.PENTAGON,
        'hexagon': MSO_SHAPE.HEXAGON,
        'octagon': MSO_SHAPE.OCTAGON,
        'star': MSO_SHAPE.STAR_5,
        'arrow_right': MSO_SHAPE.RIGHT_ARROW,
        'arrow_left': MSO_SHAPE.LEFT_ARROW,
        'arrow_up': MSO_SHAPE.UP_ARROW,
        'arrow_down': MSO_SHAPE.DOWN_ARROW,
        'callout': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        'cloud': MSO_SHAPE.CLOUD_CALLOUT,
    }

    @staticmethod
    def add_shape_slide(handler, title: str, shape_type: str,
                       text: str = "", left: float = 3.0, top: float = 2.5,
                       width: float = 4.0, height: float = 3.0,
                       fill_color: Optional[Tuple[int, int, int]] = None,
                       line_color: Optional[Tuple[int, int, int]] = None,
                       layout_index: int = 1):
        """
        Add a slide with a single shape.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            shape_type: Type of shape (see SHAPES dict)
            text: Text to add to shape
            left, top: Position in inches
            width, height: Size in inches
            fill_color: RGB tuple for fill
            line_color: RGB tuple for line
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Get shape type
        shape_enum = ShapesBuilder.SHAPES.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)

        # Add shape
        shape = slide.shapes.add_shape(
            shape_enum,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        # Apply styling
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(68, 114, 196)

        if line_color:
            shape.line.color.rgb = RGBColor(*line_color)
        else:
            shape.line.color.rgb = RGBColor(255, 255, 255)

        shape.line.width = Pt(2)

        # Add text if provided
        if text and hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            text_frame.text = text
            text_frame.word_wrap = True

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True

        return slide

    @staticmethod
    def add_callout_slide(handler, title: str, callouts: List[Dict],
                         layout_index: int = 1):
        """
        Add a slide with multiple callouts/annotations.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            callouts: List of dicts with 'text', 'left', 'top', 'width', 'height'
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        for callout in callouts:
            text = callout.get('text', '')
            left = callout.get('left', 1.0)
            top = callout.get('top', 2.0)
            width = callout.get('width', 2.5)
            height = callout.get('height', 1.0)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )

            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 200)
            shape.line.color.rgb = RGBColor(100, 100, 100)

            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                text_frame.text = text
                text_frame.word_wrap = True
                text_frame.paragraphs[0].font.size = Pt(12)

        return slide

    @staticmethod
    def add_flowchart_slide(handler, title: str, steps: List[Dict],
                           layout_index: int = 1):
        """
        Add a flowchart slide.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            steps: List of dicts with 'text', 'shape', 'decision' (bool)
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_steps = len(steps)
        box_height = 4.5 / num_steps if num_steps > 0 else 1.0
        current_top = 2.0

        for i, step in enumerate(steps):
            text = step.get('text', '')
            is_decision = step.get('decision', False)

            # Choose shape based on whether it's a decision
            if is_decision:
                shape_type = MSO_SHAPE.DIAMOND
                width, height = 3.0, 0.8
            else:
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                width, height = 3.5, 0.7

            left = 3.25  # Center

            # Add shape
            shape = slide.shapes.add_shape(
                shape_type,
                Inches(left), Inches(current_top),
                Inches(width), Inches(height)
            )

            # Style
            if is_decision:
                color = (237, 125, 49)  # Orange for decisions
            else:
                color = (68, 114, 196)  # Blue for processes

            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*color)
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)

            # Text
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                text_frame.text = text
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_SHAPE.MIDDLE

                paragraph = text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True

            # Add arrow to next step
            if i < num_steps - 1:
                arrow_top = current_top + height
                arrow_height = (box_height - height) * 0.6

                arrow = slide.shapes.add_connector(
                    1,  # Straight
                    Inches(left + width / 2), Inches(arrow_top),
                    Inches(left + width / 2), Inches(arrow_top + arrow_height)
                )
                arrow.line.color.rgb = RGBColor(100, 100, 100)
                arrow.line.width = Pt(2)

            current_top += box_height

        return slide

    @staticmethod
    def add_icon_grid_slide(handler, title: str, icons: List[Dict],
                           cols: int = 3, layout_index: int = 1):
        """
        Add a slide with a grid of icons (shapes) with labels.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            icons: List of dicts with 'shape', 'label', 'color'
            cols: Number of columns
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        num_icons = len(icons)
        rows = (num_icons + cols - 1) // cols

        icon_size = 0.8
        spacing_h = 8.0 / cols
        spacing_v = 4.5 / rows if rows > 0 else 1.5

        start_left = 1.0
        start_top = 2.0

        for i, icon in enumerate(icons):
            row = i // cols
            col = i % cols

            left = start_left + col * spacing_h
            top = start_top + row * spacing_v

            # Get shape and color
            shape_type = icon.get('shape', 'circle')
            shape_enum = ShapesBuilder.SHAPES.get(shape_type.lower(), MSO_SHAPE.OVAL)
            color = icon.get('color', (68, 114, 196))
            label = icon.get('label', '')

            # Add shape
            shape = slide.shapes.add_shape(
                shape_enum,
                Inches(left + spacing_h/2 - icon_size/2), Inches(top),
                Inches(icon_size), Inches(icon_size)
            )

            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*color)
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)

            # Add label below
            if label:
                handler.add_text_box(
                    slide, label,
                    left, top + icon_size + 0.1,
                    spacing_h, 0.4,
                    font_size=11
                )

        return slide

    @staticmethod
    def add_connector(slide, start_x: float, start_y: float,
                     end_x: float, end_y: float,
                     connector_type: int = 1,
                     color: Tuple[int, int, int] = (100, 100, 100)):
        """
        Add a connector line between two points.

        Args:
            slide: Slide object
            start_x, start_y: Start position in inches
            end_x, end_y: End position in inches
            connector_type: 1=straight, 2=elbow, 3=curved
            color: RGB tuple

        Returns:
            The connector shape
        """
        connector = slide.shapes.add_connector(
            connector_type,
            Inches(start_x), Inches(start_y),
            Inches(end_x), Inches(end_y)
        )

        connector.line.color.rgb = RGBColor(*color)
        connector.line.width = Pt(2)

        return connector

    @staticmethod
    def add_annotation_slide(handler, title: str,
                            main_content: str,
                            annotations: List[Dict],
                            layout_index: int = 1):
        """
        Add a slide with main content and annotations pointing to it.

        Args:
            handler: PPTXHandler instance
            title: Slide title
            main_content: Main text content
            annotations: List of dicts with 'text' and 'position' ('top-left', etc.)
            layout_index: Layout to use

        Returns:
            The created slide
        """
        slide = handler.add_slide(layout_index)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # Main content box
        main_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(2.5),
            Inches(5.0), Inches(3.0)
        )

        main_box.fill.solid()
        main_box.fill.fore_color.rgb = RGBColor(217, 225, 242)
        main_box.line.color.rgb = RGBColor(68, 114, 196)
        main_box.line.width = Pt(2)

        if hasattr(main_box, 'text_frame'):
            text_frame = main_box.text_frame
            text_frame.text = main_content
            text_frame.word_wrap = True
            text_frame.paragraphs[0].font.size = Pt(14)

        # Add annotations
        positions = {
            'top-left': (0.5, 2.0),
            'top-right': (8.0, 2.0),
            'bottom-left': (0.5, 4.5),
            'bottom-right': (8.0, 4.5)
        }

        for annotation in annotations:
            text = annotation.get('text', '')
            position = annotation.get('position', 'top-left')
            left, top = positions.get(position, (0.5, 2.0))

            callout = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
                Inches(left), Inches(top),
                Inches(2.0), Inches(0.8)
            )

            callout.fill.solid()
            callout.fill.fore_color.rgb = RGBColor(255, 255, 200)
            callout.line.color.rgb = RGBColor(100, 100, 100)

            if hasattr(callout, 'text_frame'):
                text_frame = callout.text_frame
                text_frame.text = text
                text_frame.word_wrap = True
                text_frame.paragraphs[0].font.size = Pt(10)

        return slide
